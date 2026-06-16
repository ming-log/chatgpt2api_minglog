from __future__ import annotations

import base64
import io
import json
import re
import threading
import time
import uuid
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections.abc import Callable, Iterable
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import unquote_to_bytes, urljoin
from urllib.request import Request, urlopen
from xml.sax.saxutils import escape

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Emu

from services.account_service import normalize_image_account_category
from services.config import DATA_DIR, config
from services.protocol import openai_v1_chat_complete, openai_v1_image_edit, openai_v1_image_generations, openai_v1_models

DEFAULT_SLIDE_COUNT = 20
MIN_SLIDE_COUNT = 1
MAX_SLIDE_COUNT = 100
AUTO_SLIDE_COUNT = "auto"
DEFAULT_IMAGE_CONCURRENCY = 10
MIN_IMAGE_CONCURRENCY = 1
MAX_IMAGE_CONCURRENCY = 100
DEFAULT_TEXT_MODEL = "gpt-5.5"
DEFAULT_IMAGE_MODEL = "gpt-image-2"
DEFAULT_IMAGE_ACCOUNT_TYPE = "free"
DEFAULT_IMAGE_QUALITY = "auto"
DEFAULT_PPT_IMAGE_SIZE = "1024x576"
PPTX_BUILD_VERSION = 4

TASK_TYPE_CONTENT = "content"
TASK_TYPE_MASTER = "master"
TASK_TYPE_PLAN = "plan"
TASK_TYPES = {TASK_TYPE_CONTENT, TASK_TYPE_MASTER, TASK_TYPE_PLAN}

IMAGE_ACCOUNT_TYPES = {"free", "paid"}
IMAGE_QUALITIES = {"auto", "low", "medium", "high"}
PPT_IMAGE_SIZES = {
    "1k": "1024x576",
    "2k": "2048x1152",
    "4k": "3840x2160",
    "1024x576": "1024x576",
    "2048x1152": "2048x1152",
    "3840x2160": "3840x2160",
}

MASTER_LAYOUT_COVER = "cover"
MASTER_LAYOUT_AGENDA = "agenda"
MASTER_LAYOUT_SECTION_BREAK = "section_break"
MASTER_LAYOUT_SINGLE_COLUMN = "single_column"
MASTER_LAYOUT_TWO_COLUMN = "two_column"
MASTER_LAYOUT_BENTO_CARD = "bento_card"
MASTER_LAYOUT_DASHBOARD = "dashboard"
MASTER_LAYOUT_THANK_YOU = "thank_you"
MASTER_BASE_LAYOUT = MASTER_LAYOUT_COVER
MASTER_LAYOUTS = {
    MASTER_LAYOUT_COVER,
    MASTER_LAYOUT_AGENDA,
    MASTER_LAYOUT_SECTION_BREAK,
    MASTER_LAYOUT_SINGLE_COLUMN,
    MASTER_LAYOUT_TWO_COLUMN,
    MASTER_LAYOUT_BENTO_CARD,
    MASTER_LAYOUT_DASHBOARD,
    MASTER_LAYOUT_THANK_YOU,
}
MASTER_LAYOUT_ORDER = [
    MASTER_LAYOUT_COVER,
    MASTER_LAYOUT_AGENDA,
    MASTER_LAYOUT_SECTION_BREAK,
    MASTER_LAYOUT_SINGLE_COLUMN,
    MASTER_LAYOUT_TWO_COLUMN,
    MASTER_LAYOUT_BENTO_CARD,
    MASTER_LAYOUT_DASHBOARD,
    MASTER_LAYOUT_THANK_YOU,
]
MASTER_LAYOUT_LABELS = {
    MASTER_LAYOUT_COVER: "封面页（Cover）",
    MASTER_LAYOUT_AGENDA: "目录页（Agenda）",
    MASTER_LAYOUT_SECTION_BREAK: "章节过渡页（Section Break）",
    MASTER_LAYOUT_SINGLE_COLUMN: "单栏内容页（Single Column）",
    MASTER_LAYOUT_TWO_COLUMN: "双栏图文页（Two-column）",
    MASTER_LAYOUT_BENTO_CARD: "卡片布局页（Bento / Card）",
    MASTER_LAYOUT_DASHBOARD: "数据图表页（Dashboard）",
    MASTER_LAYOUT_THANK_YOU: "结束页（Thank You / Q&A）",
}

TASK_STATUS_DRAFT = "draft"
TASK_STATUS_QUEUED = "queued"
TASK_STATUS_RUNNING = "running"
TASK_STATUS_SUCCESS = "success"
TASK_STATUS_ERROR = "error"
TASK_STATUS_STOPPED = "stopped"
TASK_STATUS_PACKAGING = "packaging"
TASK_STATUS_PACKAGED = "packaged"
TASK_STATUSES = {
    TASK_STATUS_DRAFT,
    TASK_STATUS_QUEUED,
    TASK_STATUS_RUNNING,
    TASK_STATUS_SUCCESS,
    TASK_STATUS_ERROR,
    TASK_STATUS_STOPPED,
    TASK_STATUS_PACKAGING,
    TASK_STATUS_PACKAGED,
}
UNFINISHED_TASK_STATUSES = {TASK_STATUS_QUEUED, TASK_STATUS_RUNNING, TASK_STATUS_PACKAGING}

SLIDE_STATUS_DRAFT = "draft"
SLIDE_STATUS_QUEUED = "queued"
SLIDE_STATUS_RUNNING = "running"
SLIDE_STATUS_SUCCESS = "success"
SLIDE_STATUS_ERROR = "error"
SLIDE_STATUS_STOPPED = "stopped"
SLIDE_STATUSES = {SLIDE_STATUS_DRAFT, SLIDE_STATUS_QUEUED, SLIDE_STATUS_RUNNING, SLIDE_STATUS_SUCCESS, SLIDE_STATUS_ERROR, SLIDE_STATUS_STOPPED}
UNFINISHED_SLIDE_STATUSES = {SLIDE_STATUS_QUEUED, SLIDE_STATUS_RUNNING}

NO_PAGE_NUMBER_PROMPT = (
    "禁止显示页码：任何页面都不得出现页码、页脚页码、幻灯片编号、总页数、"
    "Page 1、1/10、01/20 等任何表示页序的数字或文字。"
)
FONT_STYLE_PROMPT = "字体规范：中文固定使用微软雅黑（Microsoft YaHei），英文固定使用 Times New Roman。"
TITLE_STYLE_PROMPT = (
    "标题规范：同一套 PPT 内相同版式的标题必须保持固定位置、字号、字重和行高；"
    "内容页一级标题统一为约 36pt（允许 34-40pt 内微调但整套一致），二级标题约 24pt，正文约 18pt；"
    "标题过长时优先换行或压缩文案，不得随机缩放标题字号。"
)
GLOBAL_STYLE_PROMPT = (
    "全局设计规范：制作一套有高级感、可直接用于评审汇报和业务路演的 16:9 PPT 页面。"
    "优先使用深色、浓郁中性色或有质感的品牌色背景，搭配 1-2 个高识别强调色，"
    "通过大留白、清晰层级、模块化栅格、克制材质和明确内容区域建立统一视觉语言。"
    "中文固定使用微软雅黑（Microsoft YaHei），英文固定使用 Times New Roman，"
    "所有页面应保持一致的品牌调性、标题位置、图表风格和注记尺度，"
    f"{TITLE_STYLE_PROMPT}"
    f"{NO_PAGE_NUMBER_PROMPT}"
    "避免黑白单调、浅黄色泛底、过亮浅色调、随机细线、无意义网格、杂乱装饰、拥挤排版和难以阅读的小字。"
)
DEFAULT_MASTER_STYLE_PROMPT = (
    "高端商务科技风：以深色或浓郁中性色为主背景，辅以少量品牌强调色、轻微材质层次和稳定网格；"
    "版式干净、有留白、有标题舞台感，禁止浅黄色泛底、黑白单调和无关线条。"
)
MASTER_TEMPLATE_BASE_PROMPT = (
    "生成一张 16:9 PPT 母版图片，作为后续内容页图生图参考。"
    f"{FONT_STYLE_PROMPT}"
    f"{TITLE_STYLE_PROMPT}"
    "母版必须具备明确的可复用版式、统一的背景语言、安全边距和内容区域。"
    "这些规范是隐藏生成约束，不要把本段规范、字体名称或说明文字写进画面。"
)
FINAL_IMAGE_PROMPT_SUFFIX = (
    "输出完整 16:9 PPT 页面图片。画面中只呈现最终幻灯片，不要出现软件界面、边框外画布、"
    "说明文字或多张缩略图；所有文字、图表和装饰元素都必须清晰且位于安全边距内。"
    f"{NO_PAGE_NUMBER_PROMPT}"
)
PPT_STRUCTURE_PROMPT = (
    "PPT 必须按固定结构组织：第 1 页为封面页（cover），保留大留白并突出标题区域；"
    "第 2 页为目录页（agenda），强调层级感、对齐和章节编号位置；"
    "第 3 页开始按章节循环，每个章节先出现章节过渡页（section_break），再跟随若干内容页；"
    "内容页按信息类型选择 single_column、two_column、bento_card 或 dashboard；"
    "页数允许时最后一页使用 thank_you，保持极简收束，不包含联系电话、邮箱、二维码或任何联系方式。"
)

PRESENTATION_WIDTH = 12192000
PRESENTATION_HEIGHT = 6858000


class PptTaskNotFoundError(KeyError):
    pass


class PptPlanParseError(ValueError):
    def __init__(
        self,
        message: str,
        *,
        raw_output: str = "",
        attempts: list[dict[str, Any]] | None = None,
        repair_error: dict[str, Any] | None = None,
        context_exhausted: bool = False,
    ):
        super().__init__(message)
        self.raw_output = raw_output
        self.attempts = attempts or []
        self.repair_error = repair_error
        self.context_exhausted = context_exhausted

    def to_detail(self) -> dict[str, Any]:
        last_attempt = self.attempts[-1] if self.attempts else {}
        if self.context_exhausted:
            message = "当前任务过于复杂，内置模型上下文长度已耗尽，方案生成失败。请尝试配置外部文本服务后重试。"
            detail: dict[str, Any] = {
                "error": message,
                "error_type": "ppt_plan_context_exhausted",
                "hint": message,
                "model_output_preview": _text_preview(self.raw_output, 2400),
            }
            if last_attempt:
                detail["parse_error"] = last_attempt.get("message")
                detail["line"] = last_attempt.get("line")
                detail["column"] = last_attempt.get("column")
                detail["char"] = last_attempt.get("char")
                detail["error_excerpt"] = last_attempt.get("excerpt")
            if self.repair_error:
                detail["repair_error"] = self.repair_error
            return detail
        detail: dict[str, Any] = {
            "error": str(self),
            "error_type": "ppt_plan_parse_error",
            "hint": "文本模型返回的设计方案不是合法 JSON。后端已尝试自动修复；仍失败时请重试，或减少页数/简化 Markdown。",
            "model_output_preview": _text_preview(self.raw_output, 2400),
        }
        if last_attempt:
            detail["parse_error"] = last_attempt.get("message")
            detail["line"] = last_attempt.get("line")
            detail["column"] = last_attempt.get("column")
            detail["char"] = last_attempt.get("char")
            detail["error_excerpt"] = last_attempt.get("excerpt")
        if self.repair_error:
            detail["repair_error"] = self.repair_error
        return detail


def _now_iso() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _clean(value: object, default: str = "") -> str:
    return str(value or default).strip()


def _text_preview(value: object, limit: int = 1000) -> str:
    text = str(value or "")
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "..."


def _owner_id(identity: dict[str, object]) -> str:
    return _clean(identity.get("id")) or "anonymous"


def _task_key(owner_id: str, task_id: str) -> str:
    return f"{owner_id}:{task_id}"


def _safe_file_stem(value: str, fallback: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9._-]+", "-", value.strip()).strip("-._")
    return (stem or fallback)[:80]


def _pptx_build_version(value: object) -> int:
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def normalize_slide_count(value: object = None) -> int:
    if value is None or value == "":
        count = DEFAULT_SLIDE_COUNT
    elif isinstance(value, bool):
        raise ValueError(f"slide_count must be an integer between {MIN_SLIDE_COUNT} and {MAX_SLIDE_COUNT}")
    else:
        try:
            count = int(value)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"slide_count must be an integer between {MIN_SLIDE_COUNT} and {MAX_SLIDE_COUNT}") from exc
    if count < MIN_SLIDE_COUNT or count > MAX_SLIDE_COUNT:
        raise ValueError(f"slide_count must be an integer between {MIN_SLIDE_COUNT} and {MAX_SLIDE_COUNT}")
    return count


def is_auto_slide_count(value: object) -> bool:
    return isinstance(value, str) and value.strip().lower() == AUTO_SLIDE_COUNT


def normalize_image_concurrency(value: object = None) -> int:
    if value is None or value == "":
        concurrency = DEFAULT_IMAGE_CONCURRENCY
    elif isinstance(value, bool):
        raise ValueError(
            f"concurrency must be an integer between {MIN_IMAGE_CONCURRENCY} and {MAX_IMAGE_CONCURRENCY}"
        )
    else:
        try:
            concurrency = int(value)
        except (TypeError, ValueError) as exc:
            raise ValueError(
                f"concurrency must be an integer between {MIN_IMAGE_CONCURRENCY} and {MAX_IMAGE_CONCURRENCY}"
            ) from exc
    if concurrency < MIN_IMAGE_CONCURRENCY or concurrency > MAX_IMAGE_CONCURRENCY:
        raise ValueError(
            f"concurrency must be an integer between {MIN_IMAGE_CONCURRENCY} and {MAX_IMAGE_CONCURRENCY}"
        )
    return concurrency


def normalize_image_account_type(value: object = None) -> str:
    normalized = normalize_image_account_category(value)
    return normalized if normalized in IMAGE_ACCOUNT_TYPES else DEFAULT_IMAGE_ACCOUNT_TYPE


def normalize_image_quality(value: object = None) -> str:
    normalized = _clean(value, DEFAULT_IMAGE_QUALITY).lower()
    return normalized if normalized in IMAGE_QUALITIES else DEFAULT_IMAGE_QUALITY


def normalize_ppt_image_size(value: object = None) -> str:
    normalized = _clean(value).lower().replace(" ", "")
    return PPT_IMAGE_SIZES.get(normalized, DEFAULT_PPT_IMAGE_SIZE)


def _normalize_master_layout(value: object, default: str = MASTER_LAYOUT_SINGLE_COLUMN) -> str:
    clean = _clean(value).lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "pure_background": MASTER_LAYOUT_COVER,
        "blank_background": MASTER_LAYOUT_COVER,
        "background": MASTER_LAYOUT_COVER,
        "cover_page": MASTER_LAYOUT_COVER,
        "封面": MASTER_LAYOUT_COVER,
        "封面页": MASTER_LAYOUT_COVER,
        "纯背景图": MASTER_LAYOUT_COVER,
        "纯背景": MASTER_LAYOUT_COVER,
        "agenda_page": MASTER_LAYOUT_AGENDA,
        "toc": MASTER_LAYOUT_AGENDA,
        "table_of_contents": MASTER_LAYOUT_AGENDA,
        "目录": MASTER_LAYOUT_AGENDA,
        "目录页": MASTER_LAYOUT_AGENDA,
        "section": MASTER_LAYOUT_SECTION_BREAK,
        "section_transition": MASTER_LAYOUT_SECTION_BREAK,
        "section_break": MASTER_LAYOUT_SECTION_BREAK,
        "chapter_transition": MASTER_LAYOUT_SECTION_BREAK,
        "章节过渡图": MASTER_LAYOUT_SECTION_BREAK,
        "章节过度图": MASTER_LAYOUT_SECTION_BREAK,
        "章节过渡": MASTER_LAYOUT_SECTION_BREAK,
        "章节过渡页": MASTER_LAYOUT_SECTION_BREAK,
        "single": MASTER_LAYOUT_SINGLE_COLUMN,
        "single_column": MASTER_LAYOUT_SINGLE_COLUMN,
        "title_content": MASTER_LAYOUT_SINGLE_COLUMN,
        "title_body": MASTER_LAYOUT_SINGLE_COLUMN,
        "一级标题_内容": MASTER_LAYOUT_SINGLE_COLUMN,
        "一级标题+内容": MASTER_LAYOUT_SINGLE_COLUMN,
        "单栏": MASTER_LAYOUT_SINGLE_COLUMN,
        "单栏内容页": MASTER_LAYOUT_SINGLE_COLUMN,
        "two_column": MASTER_LAYOUT_TWO_COLUMN,
        "two_columns": MASTER_LAYOUT_TWO_COLUMN,
        "image_text": MASTER_LAYOUT_TWO_COLUMN,
        "title_subtitle_content": MASTER_LAYOUT_TWO_COLUMN,
        "title_subtitle_body": MASTER_LAYOUT_TWO_COLUMN,
        "一级标题_二级标题_内容": MASTER_LAYOUT_TWO_COLUMN,
        "一级标题+二级标题+内容": MASTER_LAYOUT_TWO_COLUMN,
        "双栏": MASTER_LAYOUT_TWO_COLUMN,
        "双栏图文页": MASTER_LAYOUT_TWO_COLUMN,
        "bento": MASTER_LAYOUT_BENTO_CARD,
        "card": MASTER_LAYOUT_BENTO_CARD,
        "cards": MASTER_LAYOUT_BENTO_CARD,
        "bento_card": MASTER_LAYOUT_BENTO_CARD,
        "卡片": MASTER_LAYOUT_BENTO_CARD,
        "卡片布局页": MASTER_LAYOUT_BENTO_CARD,
        "dashboard": MASTER_LAYOUT_DASHBOARD,
        "data": MASTER_LAYOUT_DASHBOARD,
        "chart": MASTER_LAYOUT_DASHBOARD,
        "数据图表页": MASTER_LAYOUT_DASHBOARD,
        "仪表盘": MASTER_LAYOUT_DASHBOARD,
        "ending": MASTER_LAYOUT_THANK_YOU,
        "end": MASTER_LAYOUT_THANK_YOU,
        "thanks": MASTER_LAYOUT_THANK_YOU,
        "thank_you": MASTER_LAYOUT_THANK_YOU,
        "qa": MASTER_LAYOUT_THANK_YOU,
        "q&a": MASTER_LAYOUT_THANK_YOU,
        "结束页": MASTER_LAYOUT_THANK_YOU,
    }
    if clean in MASTER_LAYOUTS:
        return clean
    return aliases.get(clean, default)


def _master_layout_label(layout_type: object) -> str:
    return MASTER_LAYOUT_LABELS.get(_normalize_master_layout(layout_type), MASTER_LAYOUT_LABELS[MASTER_LAYOUT_SINGLE_COLUMN])


def _master_layout_instruction(layout_type: object) -> str:
    layout = _normalize_master_layout(layout_type)
    if layout == MASTER_LAYOUT_COVER:
        return "该页使用封面页母版：画面不要太满，保留大留白，标题区域必须明显，方便后期改字。"
    if layout == MASTER_LAYOUT_AGENDA:
        return "该页使用目录页母版：强调层级感和对齐，留出编号、章节名称和短说明位置。"
    if layout == MASTER_LAYOUT_SECTION_BREAK:
        return "该页使用章节过渡页母版：视觉冲击强但元素少，大标题居中或突出，章节编号只表示章节；同类章节页编号和标题字号保持一致。"
    if layout == MASTER_LAYOUT_TWO_COLUMN:
        return "该页使用双栏图文页母版：左右比例平衡，图片区和文字区边界明确，一级标题位置和字号固定。"
    if layout == MASTER_LAYOUT_BENTO_CARD:
        return "该页使用卡片布局页母版：卡片大小有节奏变化，并保持统一圆角、阴影和间距，一级标题位置和字号固定。"
    if layout == MASTER_LAYOUT_DASHBOARD:
        return "该页使用数据图表页母版：图表区域保持留空，不画死数据，KPI 数字区域要大，一级标题位置和字号固定。"
    if layout == MASTER_LAYOUT_THANK_YOU:
        return "该页使用结束页母版：尽量简洁，只保留感谢语或 Q&A 区域，不预留联系电话、邮箱或二维码。"
    return "该页使用单栏内容页母版：内容区不要复杂，保持统一边距和充足留白，一级标题位置和字号固定。"


def _append_prompt_instruction(prompt: str, instruction: str) -> str:
    clean_prompt = _clean(prompt)
    clean_instruction = _clean(instruction)
    if not clean_instruction or clean_instruction in clean_prompt:
        return clean_prompt
    if not clean_prompt:
        return clean_instruction
    return f"{clean_prompt}\n{clean_instruction}"


def _strip_generated_chapter_instructions(prompt: str) -> str:
    lines = _clean(prompt).splitlines()
    kept: list[str] = []
    skipping_agenda_block = False
    generated_prefixes = (
        "禁止显示页码：",
        "封面页（Cover）：",
        "目录页（Agenda）：",
        "第一章章节过渡页（Section Break）：",
        "结束页（Thank You / Q&A）：",
    )
    for line in lines:
        clean_line = line.strip()
        if clean_line.startswith("目录章节清单（唯一编号标准）："):
            skipping_agenda_block = True
            continue
        if skipping_agenda_block:
            if clean_line.startswith("目录页必须按上述顺序和编号展示章节"):
                skipping_agenda_block = False
            continue
        if clean_line.startswith("章节编号一致性："):
            continue
        if clean_line.startswith("内容页章节归属："):
            continue
        if clean_line.startswith(generated_prefixes):
            continue
        kept.append(line)
    return "\n".join(kept).strip()


def _normalize_chapter_no(value: object, index: int) -> str:
    clean = _clean(value)
    if not clean:
        return f"{index:02d}"
    if re.fullmatch(r"\d+", clean):
        return f"{int(clean):02d}"
    return clean


def _strip_chapter_heading(value: object) -> str:
    text = _clean(value)
    if not text:
        return ""
    patterns = [
        r"^第\s*[一二三四五六七八九十百千万零〇两\d]+\s*[章节篇部分]\s*[：:、.．\-\s]*",
        r"^(?:章节|章|部分)\s*[一二三四五六七八九十百千万零〇两\d]+\s*[：:、.．\-\s]*",
        r"^(?:chapter|section|part)\s*[A-Za-z0-9IVXLCDM]+\s*[：:、.．\-\s]*",
        r"^\d{1,3}\s*[：:、.．\-\s]+",
        r"^[一二三四五六七八九十百千万零〇两]+\s*[、.．]\s*",
    ]
    for pattern in patterns:
        text = re.sub(pattern, "", text, count=1, flags=re.IGNORECASE).strip()
    return text


def _looks_generic_chapter_title(value: object) -> bool:
    clean = _clean(value).lower()
    if not clean:
        return True
    if re.fullmatch(r"第\s*\d+\s*页", clean):
        return True
    return clean in {
        "章节",
        "章节页",
        "章节过渡",
        "章节过渡页",
        "section",
        "section break",
        "chapter",
        "chapter transition",
    }


def _chapter_title_from_slide(slide: dict[str, Any], index: int) -> str:
    for key in ("chapter_title", "section_title", "chapter_name", "section_name"):
        candidate = _strip_chapter_heading(slide.get(key))
        if not _looks_generic_chapter_title(candidate):
            return candidate[:80]
    title_candidate = _strip_chapter_heading(slide.get("title"))
    if not _looks_generic_chapter_title(title_candidate):
        return title_candidate[:80]
    prompt = _strip_generated_chapter_instructions(_clean(slide.get("slide_prompt")))
    for line in prompt.splitlines():
        candidate = _strip_chapter_heading(line)
        if not _looks_generic_chapter_title(candidate):
            return candidate[:80]
    return f"第 {index} 章"


def _strip_prompt_leading_chapter_no(prompt: str) -> str:
    lines = _clean(prompt).splitlines()
    if not lines:
        return ""
    first = _strip_chapter_heading(lines[0])
    if first and first != lines[0].strip():
        lines[0] = first
    return "\n".join(lines).strip()


def _normalize_plan_chapters(value: object) -> list[dict[str, str]]:
    if not isinstance(value, list):
        return []
    chapters: list[dict[str, str]] = []
    for index, item in enumerate(value, start=1):
        if isinstance(item, dict):
            title = _clean(
                item.get("title")
                or item.get("chapter_title")
                or item.get("section_title")
                or item.get("name")
                or item.get("chapter_name")
                or item.get("section_name")
            )
            chapter_no = item.get("chapter_no") or item.get("chapter_number") or item.get("section_no") or item.get("number")
        else:
            title = _clean(item)
            chapter_no = ""
        title = _strip_chapter_heading(title) or f"第 {index} 章"
        chapter = {
            "chapter_no": _normalize_chapter_no(chapter_no, index),
            "chapter_title": title[:80],
        }
        if isinstance(item, dict) and _clean(item.get("section_slide_id")):
            chapter["section_slide_id"] = _clean(item.get("section_slide_id"))
        chapters.append(chapter)
    return chapters


def _apply_chapter_consistency(
    slides: list[dict[str, Any]],
    plan_chapters: list[dict[str, str]] | None = None,
) -> list[dict[str, str]]:
    section_positions = [
        index
        for index, slide in enumerate(slides)
        if _normalize_master_layout(slide.get("layout_type")) == MASTER_LAYOUT_SECTION_BREAK
    ]
    if not section_positions:
        return []
    chapter_specs = plan_chapters or []
    chapters: list[dict[str, str]] = []
    for chapter_index, position in enumerate(section_positions, start=1):
        slide = slides[position]
        spec = chapter_specs[chapter_index - 1] if chapter_index - 1 < len(chapter_specs) else {}
        chapter_no = _normalize_chapter_no(spec.get("chapter_no"), chapter_index)
        derived_title = _chapter_title_from_slide(slide, chapter_index)
        spec_title = _clean(spec.get("chapter_title"))
        chapter_title = spec_title if derived_title == f"第 {chapter_index} 章" and spec_title else derived_title
        slide["chapter_no"] = chapter_no
        slide["chapter_title"] = chapter_title
        chapters.append({
            "chapter_no": chapter_no,
            "chapter_title": chapter_title,
            "section_slide_id": _clean(slide.get("slide_id"), str(position + 1)),
        })
        next_position = section_positions[chapter_index] if chapter_index < len(section_positions) else len(slides)
        for content_slide in slides[position:next_position]:
            layout = _normalize_master_layout(content_slide.get("layout_type"))
            if layout in {MASTER_LAYOUT_COVER, MASTER_LAYOUT_AGENDA, MASTER_LAYOUT_THANK_YOU}:
                continue
            content_slide["chapter_no"] = chapter_no
            content_slide["chapter_title"] = chapter_title
        slide["slide_prompt"] = _strip_prompt_leading_chapter_no(
            _strip_generated_chapter_instructions(_clean(slide.get("slide_prompt")))
        )
    return chapters


def _task_slides(task: dict[str, Any]) -> list[dict[str, Any]]:
    slides = task.get("slides")
    return slides if isinstance(slides, list) else []


def _chapters_from_task(task: dict[str, Any]) -> list[dict[str, str]]:
    chapters = _normalize_plan_chapters(task.get("chapters"))
    if chapters:
        return chapters
    derived: list[dict[str, str]] = []
    for index, slide in enumerate(
        [item for item in _task_slides(task) if _normalize_master_layout(item.get("layout_type")) == MASTER_LAYOUT_SECTION_BREAK],
        start=1,
    ):
        derived.append({
            "chapter_no": _normalize_chapter_no(slide.get("chapter_no"), index),
            "chapter_title": _clean(slide.get("chapter_title")) or _chapter_title_from_slide(slide, index),
            "section_slide_id": _clean(slide.get("slide_id")),
        })
    return derived


def _chapter_for_slide(task: dict[str, Any], slide: dict[str, Any]) -> tuple[int, dict[str, str]]:
    chapters = _chapters_from_task(task)
    slide_id = _clean(slide.get("slide_id"))
    slide_chapter_no = _clean(slide.get("chapter_no"))
    slide_chapter_title = _clean(slide.get("chapter_title"))
    for index, chapter in enumerate(chapters, start=1):
        if _clean(chapter.get("section_slide_id")) and _clean(chapter.get("section_slide_id")) == slide_id:
            return index, chapter
        if slide_chapter_no and _clean(chapter.get("chapter_no")) == slide_chapter_no:
            return index, chapter
        if slide_chapter_title and _clean(chapter.get("chapter_title")) == slide_chapter_title:
            return index, chapter
    fallback_index = max(1, len(chapters) + 1)
    fallback = {
        "chapter_no": _normalize_chapter_no(slide.get("chapter_no"), fallback_index),
        "chapter_title": slide_chapter_title or _chapter_title_from_slide(slide, fallback_index),
        "section_slide_id": slide_id,
    }
    return fallback_index, fallback


def _presentation_theme(task: dict[str, Any], slide: dict[str, Any]) -> str:
    return (
        _clean(slide.get("title"))
        or _clean(task.get("name"))
        or _clean(task.get("markdown_file_name"))
        or _clean(task.get("design_concept"))
        or "本套 PPT 主题"
    )


def _slide_generation_control_prompt(task: dict[str, Any], slide: dict[str, Any]) -> str:
    layout = _normalize_master_layout(slide.get("layout_type"))
    chapters = _chapters_from_task(task)
    if layout == MASTER_LAYOUT_COVER:
        theme = _presentation_theme(task, slide)
        return f"封面页硬性要求：必须清晰呈现整套 PPT 的主题「{theme}」；不要出现目录列表、正文段落、页码或章节编号。"
    if layout == MASTER_LAYOUT_AGENDA:
        if chapters:
            outline = "\n".join(f"{chapter['chapter_no']} {chapter['chapter_title']}" for chapter in chapters)
            return (
                f"目录页硬性要求：必须提供 {len(chapters)} 个章节标题，并严格按以下唯一章节清单展示：\n"
                f"{outline}\n"
                "目录页只展示章节编号、章节标题和必要短说明；不得显示页码、每章页数或联系方式。"
            )
        return "目录页硬性要求：必须提供从正文内容归纳出的章节标题；不得显示页码、每章页数或联系方式。"
    if layout == MASTER_LAYOUT_SECTION_BREAK:
        chapter_index, chapter = _chapter_for_slide(task, slide)
        chapter_no = _clean(chapter.get("chapter_no"), _normalize_chapter_no(slide.get("chapter_no"), chapter_index))
        chapter_title = _clean(chapter.get("chapter_title")) or _chapter_title_from_slide(slide, chapter_index)
        return (
            f"章节过渡页硬性要求：这是第 {chapter_index} 章的章节过渡页；"
            f"当前章节编号必须显示为「{chapter_no}」，当前章节标题必须显示为「{chapter_title}」。"
            "只呈现当前章节，不得出现其它章节编号、目录列表、页码、总页数或联系方式。"
        )
    if layout == MASTER_LAYOUT_THANK_YOU:
        return (
            "结束页硬性要求：保持极简收束，只允许出现感谢语、Q&A 或简短结束标题；"
            "不允许包含联系电话、手机号、电子信箱、邮箱地址、个人微信、二维码、网址或任何联系方式。"
        )
    chapter_no = _clean(slide.get("chapter_no"))
    chapter_title = _clean(slide.get("chapter_title"))
    if chapter_no or chapter_title:
        return f"内容页章节归属：本页属于章节「{chapter_no} {chapter_title}」，不要显示页码或联系方式。"
    return ""


def _master_style_context(style_prompt: object = "") -> str:
    clean_style = _clean(style_prompt)
    style = clean_style or DEFAULT_MASTER_STYLE_PROMPT
    return f"{GLOBAL_STYLE_PROMPT}\n母版整体风格：{style}"


def _reference_image_from_slide(slide: dict[str, Any], *, fallback_id: str = "") -> dict[str, str]:
    layout_type = _normalize_master_layout(slide.get("layout_type") or slide.get("slide_id"))
    reference_id = _clean(slide.get("reference_id") or slide.get("slide_id") or fallback_id or layout_type)
    return {
        "id": reference_id,
        "title": _clean(slide.get("title"), _master_layout_label(layout_type)),
        "layout_type": layout_type,
        "image_url": _clean(slide.get("image_url")),
    }


def _normalize_reference_images(value: object) -> list[dict[str, str]]:
    if not isinstance(value, list):
        return []
    references: list[dict[str, str]] = []
    for index, item in enumerate(value, start=1):
        if not isinstance(item, dict):
            continue
        image_url = _clean(item.get("image_url"))
        if not image_url:
            continue
        layout_type = _normalize_master_layout(item.get("layout_type") or item.get("id"))
        reference_id = _clean(item.get("id") or item.get("slide_id"), f"ref-{index}")
        references.append({
            "id": reference_id,
            "title": _clean(item.get("title"), _master_layout_label(layout_type)),
            "layout_type": layout_type,
            "image_url": image_url,
        })
    return references


def _apply_required_ppt_structure(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    structured = [dict(slide) for slide in slides]
    for index, slide in enumerate(structured, start=1):
        slide["slide_id"] = str(index)
        slide["slide_prompt"] = _strip_generated_chapter_instructions(slide.get("slide_prompt", ""))

    if not structured:
        return structured

    title = _clean(structured[0].get("title"), "封面")
    structured[0]["title"] = title
    structured[0]["layout_type"] = MASTER_LAYOUT_COVER
    structured[0]["slide_prompt"] = _strip_generated_chapter_instructions(structured[0].get("slide_prompt", ""))

    if len(structured) >= 2:
        structured[1]["title"] = "目录"
        structured[1]["layout_type"] = MASTER_LAYOUT_AGENDA
        structured[1]["slide_prompt"] = _strip_generated_chapter_instructions(structured[1].get("slide_prompt", ""))

    if len(structured) >= 3:
        structured[2]["layout_type"] = MASTER_LAYOUT_SECTION_BREAK
        structured[2]["slide_prompt"] = _strip_prompt_leading_chapter_no(
            _strip_generated_chapter_instructions(structured[2].get("slide_prompt", ""))
        )

    if len(structured) >= 4:
        structured[-1]["layout_type"] = MASTER_LAYOUT_THANK_YOU
        structured[-1]["slide_prompt"] = _strip_generated_chapter_instructions(structured[-1].get("slide_prompt", ""))

    return structured


def _chat_content(result: dict[str, Any]) -> str:
    choices = result.get("choices")
    if not isinstance(choices, list) or not choices:
        return ""
    first = choices[0] if isinstance(choices[0], dict) else {}
    message = first.get("message") if isinstance(first.get("message"), dict) else {}
    content = message.get("content")
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, dict):
                text = item.get("text")
                if isinstance(text, str):
                    parts.append(text)
            elif isinstance(item, str):
                parts.append(item)
        return "".join(parts).strip()
    return ""


def _chat_chunk_content(chunk: dict[str, Any]) -> str:
    choices = chunk.get("choices")
    first = choices[0] if isinstance(choices, list) and choices and isinstance(choices[0], dict) else {}
    delta = first.get("delta") if isinstance(first.get("delta"), dict) else {}
    content = delta.get("content")
    return content if isinstance(content, str) else ""


def _loads_json_object(text: str) -> dict[str, Any]:
    raw = text.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.IGNORECASE).strip()
        raw = re.sub(r"\s*```$", "", raw).strip()
    attempts: list[dict[str, Any]] = []
    candidates = [raw]
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        candidates.append(raw[start:end + 1])
    for candidate in dict.fromkeys(candidates):
        try:
            parsed = json.loads(candidate)
            break
        except json.JSONDecodeError as exc:
            attempts.append(_json_error_attempt(candidate, exc))
    else:
        if attempts:
            message = f"模型返回的设计方案不是合法 JSON：{attempts[-1].get('message')}"
        else:
            message = "模型返回的设计方案不是合法 JSON：未找到 JSON object"
        raise PptPlanParseError(message, raw_output=text, attempts=attempts)
    if not isinstance(parsed, dict):
        raise ValueError("设计方案必须是 JSON object")
    return parsed


def _json_error_attempt(text: str, exc: json.JSONDecodeError) -> dict[str, Any]:
    start = max(0, exc.pos - 240)
    end = min(len(text), exc.pos + 240)
    pointer = " " * max(0, exc.pos - start) + "^"
    return {
        "message": exc.msg,
        "line": exc.lineno,
        "column": exc.colno,
        "char": exc.pos,
        "excerpt": f"{text[start:end]}\n{pointer}",
    }


def _looks_like_truncated_json(text: str, attempts: list[dict[str, Any]] | None = None) -> bool:
    raw = text.strip()
    if not raw:
        return False
    start = raw.find("{")
    if start < 0:
        return False
    candidate = raw[start:]
    if candidate.count("{") > candidate.count("}") or candidate.count("[") > candidate.count("]"):
        return True
    last = candidate.rstrip()[-1:]
    if last and last not in {"}", "]"} and '"slides"' in candidate:
        return True
    last_attempt = (attempts or [])[-1] if attempts else {}
    message = str(last_attempt.get("message") or "").lower()
    char = last_attempt.get("char")
    if isinstance(char, int) and len(candidate) - char <= 3 and (
        "expecting" in message or "unterminated" in message
    ):
        return True
    return False


def _extract_image_url(result: dict[str, Any]) -> str:
    data = result.get("data")
    if not isinstance(data, list) or not data:
        message = _clean(result.get("message"))
        raise RuntimeError(message or "图片接口未返回图片数据")
    first = data[0] if isinstance(data[0], dict) else {}
    url = _clean(first.get("url"))
    if url:
        return url
    b64_json = _clean(first.get("b64_json"))
    if b64_json:
        return f"data:image/png;base64,{b64_json}"
    raise RuntimeError("图片接口未返回 url 或 b64_json")


def _fetch_image_bytes(url: str, base_url: str = "") -> bytes:
    value = _clean(url)
    if not value:
        raise ValueError("image_url is required")
    if value.startswith("data:"):
        header, _, payload = value.partition(",")
        if not payload:
            raise ValueError("invalid data image url")
        if ";base64" in header.lower():
            return base64.b64decode(payload)
        return unquote_to_bytes(payload)
    if value.startswith("/"):
        if not base_url:
            raise ValueError("base_url is required for relative image url")
        value = urljoin(base_url.rstrip("/") + "/", value.lstrip("/"))
    if value.startswith(("http://", "https://")):
        with urlopen(value, timeout=30) as response:
            return response.read()
    path = Path(value)
    if path.is_file():
        return path.read_bytes()
    raise ValueError(f"unsupported image url: {value}")


def _mime_from_image_bytes(data: bytes) -> str:
    try:
        with Image.open(io.BytesIO(data)) as image:
            image_format = (image.format or "").lower()
    except Exception:
        return "image/png"
    if image_format == "jpeg":
        return "image/jpeg"
    if image_format in {"png", "webp", "gif"}:
        return f"image/{image_format}"
    return "image/png"


def _extension_from_mime(mime_type: str) -> str:
    if mime_type == "image/jpeg":
        return "jpg"
    if mime_type == "image/webp":
        return "webp"
    if mime_type == "image/gif":
        return "gif"
    return "png"


def _image_metadata_from_bytes(data: bytes) -> dict[str, int]:
    metadata: dict[str, int] = {"image_size": len(data)}
    try:
        with Image.open(io.BytesIO(data)) as image:
            metadata["image_width"] = int(image.width)
            metadata["image_height"] = int(image.height)
    except Exception:
        pass
    return metadata


def _image_metadata_from_url(url: str, base_url: str = "") -> dict[str, int]:
    try:
        return _image_metadata_from_bytes(_fetch_image_bytes(url, base_url))
    except Exception:
        return {}


def _image_input_from_url(url: str, base_url: str = "") -> tuple[bytes, str, str]:
    data = _fetch_image_bytes(url, base_url)
    mime_type = _mime_from_image_bytes(data)
    extension = "jpg" if mime_type == "image/jpeg" else mime_type.split("/", 1)[-1]
    return data, f"ppt-slide.{extension}", mime_type


def _image_input_to_data_url(image: tuple[bytes, str, str]) -> str:
    data, _, mime_type = image
    return f"data:{mime_type or 'image/png'};base64,{base64.b64encode(data).decode('ascii')}"


def _openai_endpoint(base_url: str, path: str) -> str:
    root = _clean(base_url).rstrip("/")
    if not root:
        raise ValueError("OpenAI 接口地址不能为空")
    if root.endswith("/v1"):
        return f"{root}{path.removeprefix('/v1')}"
    return f"{root}{path}"


def _post_openai_json(base_url: str, api_key: str, path: str, payload: dict[str, Any]) -> dict[str, Any]:
    token = _clean(api_key)
    if not token:
        raise ValueError("使用外部 OpenAI 协议接口时必须填写 API Key")
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = Request(
        _openai_endpoint(base_url, path),
        data=data,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )
    with urlopen(request, timeout=120) as response:
        raw = response.read().decode("utf-8")
    parsed = json.loads(raw)
    if not isinstance(parsed, dict):
        raise RuntimeError("OpenAI 协议接口返回值不是 JSON object")
    return parsed


def _post_openai_stream_json(base_url: str, api_key: str, path: str, payload: dict[str, Any]) -> Iterable[dict[str, Any]]:
    token = _clean(api_key)
    if not token:
        raise ValueError("使用外部 OpenAI 协议接口时必须填写 API Key")
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = Request(
        _openai_endpoint(base_url, path),
        data=data,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "Accept": "text/event-stream",
        },
        method="POST",
    )
    with urlopen(request, timeout=120) as response:
        for raw_line in response:
            line = raw_line.decode("utf-8", errors="replace").strip()
            if not line or line.startswith(":") or not line.startswith("data:"):
                continue
            data_line = line.removeprefix("data:").strip()
            if data_line == "[DONE]":
                break
            parsed = json.loads(data_line)
            if isinstance(parsed, dict):
                yield parsed


def _get_openai_json(base_url: str, api_key: str, path: str, *, timeout: float = 30.0) -> tuple[int, dict[str, Any]]:
    token = _clean(api_key)
    if not token:
        raise ValueError("使用外部 OpenAI 协议接口时必须填写 API Key")
    request = Request(
        _openai_endpoint(base_url, path),
        headers={
            "Authorization": f"Bearer {token}",
            "Accept": "application/json",
        },
        method="GET",
    )
    with urlopen(request, timeout=timeout) as response:
        status = int(getattr(response, "status", 200) or 200)
        raw = response.read().decode("utf-8")
    parsed = json.loads(raw)
    if not isinstance(parsed, dict):
        raise RuntimeError("OpenAI 协议接口返回值不是 JSON object")
    return status, parsed


def _model_ids(result: dict[str, Any]) -> set[str]:
    data = result.get("data")
    if not isinstance(data, list):
        return set()
    ids: set[str] = set()
    for item in data:
        if not isinstance(item, dict):
            continue
        model_id = _clean(item.get("id"))
        if model_id:
            ids.add(model_id)
    return ids


def _normalize_png(image_bytes: bytes) -> bytes:
    with Image.open(io.BytesIO(image_bytes)) as image:
        image = ImageOps.exif_transpose(image)
        if image.mode not in {"RGB", "RGBA"}:
            image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
        buf = io.BytesIO()
        image.save(buf, format="PNG", optimize=True)
        return buf.getvalue()


def _xml_header() -> str:
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'


def _xml_attr(value: object) -> str:
    return escape(str(value or ""), {'"': "&quot;", "'": "&apos;"})


def _rels_xml(items: list[tuple[str, str, str]]) -> str:
    body = "".join(
        f'<Relationship Id="{_xml_attr(rel_id)}" Type="{_xml_attr(rel_type)}" Target="{_xml_attr(target)}"/>'
        for rel_id, rel_type, target in items
    )
    return f'{_xml_header()}<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">{body}</Relationships>'


def _content_types_xml(slide_count: int) -> str:
    overrides = [
        ('/docProps/core.xml', 'application/vnd.openxmlformats-package.core-properties+xml'),
        ('/docProps/app.xml', 'application/vnd.openxmlformats-officedocument.extended-properties+xml'),
        ('/ppt/presentation.xml', 'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'),
        ('/ppt/slideMasters/slideMaster1.xml', 'application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml'),
        ('/ppt/slideLayouts/slideLayout1.xml', 'application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml'),
        ('/ppt/theme/theme1.xml', 'application/vnd.openxmlformats-officedocument.theme+xml'),
    ]
    overrides.extend(
        (f'/ppt/slides/slide{index}.xml', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
        for index in range(1, slide_count + 1)
    )
    override_xml = "".join(
        f'<Override PartName="{_xml_attr(part)}" ContentType="{_xml_attr(content_type)}"/>'
        for part, content_type in overrides
    )
    return (
        f'{_xml_header()}<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Default Extension="png" ContentType="image/png"/>'
        f'{override_xml}</Types>'
    )


def _core_xml() -> str:
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    return (
        f'{_xml_header()}<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        'xmlns:dcterms="http://purl.org/dc/terms/" '
        'xmlns:dcmitype="http://purl.org/dc/dcmitype/" '
        'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        '<dc:title>Markdown generated presentation</dc:title>'
        '<dc:creator>chatgpt2api</dc:creator>'
        f'<dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>'
        f'<dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>'
        '</cp:coreProperties>'
    )


def _app_xml(slide_count: int) -> str:
    return (
        f'{_xml_header()}<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        '<Application>chatgpt2api</Application>'
        '<PresentationFormat>On-screen Show (16:9)</PresentationFormat>'
        f'<Slides>{slide_count}</Slides>'
        '</Properties>'
    )


def _presentation_xml(slide_count: int) -> str:
    slide_ids = "".join(
        f'<p:sldId id="{255 + index}" r:id="rId{index + 1}"/>'
        for index in range(1, slide_count + 1)
    )
    return (
        f'{_xml_header()}<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>'
        f'<p:sldIdLst>{slide_ids}</p:sldIdLst>'
        f'<p:sldSz cx="{PRESENTATION_WIDTH}" cy="{PRESENTATION_HEIGHT}" type="wide"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        '<p:defaultTextStyle><a:defPPr><a:defRPr lang="zh-CN"/></a:defPPr></p:defaultTextStyle>'
        '</p:presentation>'
    )


def _presentation_rels_xml(slide_count: int) -> str:
    items = [(
        "rId1",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
        "slideMasters/slideMaster1.xml",
    )]
    items.extend(
        (
            f"rId{index + 1}",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
            f"slides/slide{index}.xml",
        )
        for index in range(1, slide_count + 1)
    )
    return _rels_xml(items)


def _group_shape_xml() -> str:
    return (
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
    )


def _slide_master_xml() -> str:
    return (
        f'{_xml_header()}<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        f'{_group_shape_xml()}'
        '</p:spTree></p:cSld>'
        '<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" '
        'accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>'
        '<p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>'
        '<p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>'
        '</p:sldMaster>'
    )


def _slide_layout_xml() -> str:
    return (
        f'{_xml_header()}<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">'
        '<p:cSld name="Blank"><p:spTree>'
        f'{_group_shape_xml()}'
        '</p:spTree></p:cSld>'
        '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
        '</p:sldLayout>'
    )


def _theme_xml() -> str:
    return (
        f'{_xml_header()}<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="chatgpt2api">'
        '<a:themeElements>'
        '<a:clrScheme name="Office">'
        '<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>'
        '<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>'
        '<a:dk2><a:srgbClr val="1F2937"/></a:dk2>'
        '<a:lt2><a:srgbClr val="F8FAFC"/></a:lt2>'
        '<a:accent1><a:srgbClr val="2563EB"/></a:accent1>'
        '<a:accent2><a:srgbClr val="059669"/></a:accent2>'
        '<a:accent3><a:srgbClr val="D97706"/></a:accent3>'
        '<a:accent4><a:srgbClr val="7C3AED"/></a:accent4>'
        '<a:accent5><a:srgbClr val="DB2777"/></a:accent5>'
        '<a:accent6><a:srgbClr val="0F766E"/></a:accent6>'
        '<a:hlink><a:srgbClr val="2563EB"/></a:hlink>'
        '<a:folHlink><a:srgbClr val="7C3AED"/></a:folHlink>'
        '</a:clrScheme>'
        '<a:fontScheme name="Office"><a:majorFont><a:latin typeface="Aptos Display"/><a:ea typeface="Microsoft YaHei"/><a:cs typeface=""/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Aptos"/><a:ea typeface="Microsoft YaHei"/><a:cs typeface=""/></a:minorFont></a:fontScheme>'
        '<a:fmtScheme name="Office">'
        '<a:fillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:gradFill rotWithShape="1"><a:gsLst><a:gs pos="0"><a:schemeClr val="phClr"><a:lumMod val="110000"/><a:satMod val="105000"/></a:schemeClr></a:gs><a:gs pos="100000"><a:schemeClr val="phClr"><a:lumMod val="85000"/><a:satMod val="95000"/></a:schemeClr></a:gs></a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'
        '<a:gradFill rotWithShape="1"><a:gsLst><a:gs pos="0"><a:schemeClr val="phClr"><a:lumMod val="102000"/><a:satMod val="103000"/></a:schemeClr></a:gs><a:gs pos="100000"><a:schemeClr val="phClr"><a:lumMod val="72000"/><a:satMod val="93000"/></a:schemeClr></a:gs></a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'
        '</a:fillStyleLst>'
        '<a:lnStyleLst>'
        '<a:ln w="6350" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
        '<a:ln w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
        '<a:ln w="19050" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
        '</a:lnStyleLst>'
        '<a:effectStyleLst>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '</a:effectStyleLst>'
        '<a:bgFillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"><a:tint val="95000"/><a:satMod val="170000"/></a:schemeClr></a:solidFill>'
        '<a:gradFill rotWithShape="1"><a:gsLst><a:gs pos="0"><a:schemeClr val="phClr"><a:tint val="93000"/><a:satMod val="150000"/></a:schemeClr></a:gs><a:gs pos="100000"><a:schemeClr val="phClr"><a:shade val="98000"/><a:lumMod val="102000"/></a:schemeClr></a:gs></a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'
        '</a:bgFillStyleLst>'
        '</a:fmtScheme>'
        '</a:themeElements>'
        '</a:theme>'
    )


def _slide_xml(index: int, title: str) -> str:
    name = _xml_attr(title or f"Slide {index}")
    return (
        f'{_xml_header()}<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        f'{_group_shape_xml()}'
        '<p:pic><p:nvPicPr>'
        f'<p:cNvPr id="2" name="{name}"/>'
        '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/>'
        '</p:nvPicPr>'
        '<p:blipFill><a:blip r:embed="rId1"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{PRESENTATION_WIDTH}" cy="{PRESENTATION_HEIGHT}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        '</p:pic>'
        '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
    )


def _slide_rels_xml(index: int) -> str:
    return _rels_xml([
        (
            "rId1",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            f"../media/image{index}.png",
        ),
        (
            "rId2",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
            "../slideLayouts/slideLayout1.xml",
        ),
    ])


def build_image_pptx(slides: list[tuple[str, bytes]]) -> bytes:
    if not slides:
        raise ValueError("slides is required")
    presentation = Presentation()
    presentation.slide_width = Emu(PRESENTATION_WIDTH)
    presentation.slide_height = Emu(PRESENTATION_HEIGHT)
    presentation.core_properties.title = "Markdown generated presentation"
    presentation.core_properties.author = "chatgpt2api"
    blank_layout = presentation.slide_layouts[6]

    for title, image_bytes in slides:
        slide = presentation.slides.add_slide(blank_layout)
        picture = slide.shapes.add_picture(
            io.BytesIO(_normalize_png(image_bytes)),
            0,
            0,
            width=Emu(PRESENTATION_WIDTH),
            height=Emu(PRESENTATION_HEIGHT),
        )
        picture.name = str(title or "Slide image")[:255]

    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


class PptTaskService:
    def __init__(
        self,
        path: Path,
        *,
        package_dir: Path | None = None,
        chat_handler: Callable[[dict[str, Any]], dict[str, Any]] = openai_v1_chat_complete.handle,
        image_handler: Callable[[dict[str, Any]], dict[str, Any]] = openai_v1_image_generations.handle,
        image_edit_handler: Callable[[dict[str, Any]], dict[str, Any]] = openai_v1_image_edit.handle,
        image_fetcher: Callable[[str, str], bytes] = _fetch_image_bytes,
    ):
        self.path = path
        self.package_dir = package_dir or (DATA_DIR / "ppt_packages")
        self.chat_handler = chat_handler
        self.image_handler = image_handler
        self.image_edit_handler = image_edit_handler
        self.image_fetcher = image_fetcher
        self._lock = threading.RLock()
        self._tasks: dict[str, dict[str, Any]] = {}
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.package_dir.mkdir(parents=True, exist_ok=True)
        with self._lock:
            self._tasks = self._load_locked()
            if self._recover_unfinished_locked():
                self._save_locked()

    def create_plan(
        self,
        markdown: str,
        slide_count: object = AUTO_SLIDE_COUNT,
        *,
        model: str = DEFAULT_TEXT_MODEL,
        text_base_url: str = "",
        text_api_key: str = "",
    ) -> dict[str, Any]:
        auto_count = slide_count is None or is_auto_slide_count(slide_count)
        count = None if auto_count else normalize_slide_count(slide_count)
        source = _clean(markdown)
        if not source:
            raise ValueError("markdown is required")
        payload = {
            "model": _clean(model, DEFAULT_TEXT_MODEL),
            "stream": False,
            "messages": self._plan_messages(source, count),
        }
        content = self._chat_text(
            payload,
            stream_label=f"plan slide_count={AUTO_SLIDE_COUNT if auto_count else count}",
            base_url=text_base_url,
            api_key=text_api_key,
        )
        if not content:
            raise RuntimeError("文本接口未返回设计方案")
        try:
            return self.normalize_plan(_loads_json_object(content), count)
        except PptPlanParseError as exc:
            internal_context_exhausted = not _clean(text_base_url) and _looks_like_truncated_json(content, exc.attempts)
            return self._repair_and_normalize_plan(
                content,
                count,
                _clean(model, DEFAULT_TEXT_MODEL),
                exc,
                text_base_url=text_base_url,
                text_api_key=text_api_key,
                context_exhausted=internal_context_exhausted,
            )

    def normalize_plan(self, raw_plan: dict[str, Any], slide_count: object = None) -> dict[str, Any]:
        plan = raw_plan.get("plan") if isinstance(raw_plan.get("plan"), dict) else raw_plan
        count = normalize_slide_count(slide_count if slide_count is not None else plan.get("slide_count"))
        plan_chapters = _normalize_plan_chapters(plan.get("chapters") or plan.get("sections") or plan.get("chapter_outline"))
        slides_raw = plan.get("slides") or plan.get("pages")
        if not isinstance(slides_raw, list) or len(slides_raw) != count:
            raise ValueError("设计方案页数与 slide_count 不匹配，请重试")
        slides: list[dict[str, Any]] = []
        for index, item in enumerate(slides_raw, start=1):
            if not isinstance(item, dict):
                raise ValueError("slides 中的每一页都必须是 object")
            title = _clean(item.get("title"), f"第 {index} 页")
            slide_prompt = _strip_generated_chapter_instructions(
                _clean(item.get("slide_prompt") or item.get("prompt") or item.get("image_prompt"))
            )
            if not slide_prompt:
                raise ValueError(f"第 {index} 页缺少 slide_prompt")
            layout_type = _normalize_master_layout(
                item.get("layout_type") or item.get("master_layout") or item.get("template_type")
            )
            slides.append({
                "slide_id": _clean(item.get("slide_id") or item.get("id"), str(index)),
                "title": title,
                "layout_type": layout_type,
                "slide_prompt": slide_prompt,
                "chapter_no": _clean(item.get("chapter_no") or item.get("chapter_number") or item.get("section_no")),
                "chapter_title": _clean(item.get("chapter_title") or item.get("chapter_name") or item.get("section_title")),
            })
        slides = _apply_required_ppt_structure(slides)
        chapters = _apply_chapter_consistency(slides, plan_chapters)
        return {
            "slide_count": count,
            "design_concept": _clean(plan.get("design_concept") or plan.get("overall_design") or plan.get("summary")),
            "global_style_prompt": GLOBAL_STYLE_PROMPT,
            "chapters": chapters,
            "slides": slides,
        }

    def _repair_and_normalize_plan(
        self,
        raw_content: str,
        slide_count: int | None,
        model: str,
        original_error: PptPlanParseError,
        *,
        text_base_url: str = "",
        text_api_key: str = "",
        context_exhausted: bool = False,
    ) -> dict[str, Any]:
        payload = {
            "model": model,
            "stream": False,
            "messages": [
                {
                    "role": "system",
                    "content": (
                        "你是严格的 JSON 修复器。只输出一个合法 JSON object，不要输出 Markdown、解释或代码块。"
                        "不要新增页，不要删除页，不要改写语义，只修复引号、逗号、转义、尾逗号等 JSON 语法问题。"
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"请把下面内容修复为合法 JSON，并保持 slides 与 slide_count 精确一致。"
                        "必须包含 slide_count、design_concept、global_style_prompt、chapters、slides。"
                        "每个 slides item 必须包含 slide_id、title、layout_type、slide_prompt。"
                        f"layout_type 只能是 {', '.join(MASTER_LAYOUT_ORDER)}。\n\n"
                        f"{PPT_STRUCTURE_PROMPT}"
                        f"{NO_PAGE_NUMBER_PROMPT}\n\n"
                        f"{raw_content}"
                    ),
                },
            ],
        }
        repaired_content = self._chat_text(
            payload,
            stream_label=f"plan-repair slide_count={slide_count}",
            base_url=text_base_url,
            api_key=text_api_key,
        )
        if not repaired_content:
            original_error.repair_error = {"error": "文本接口未返回修复后的设计方案"}
            original_error.context_exhausted = context_exhausted or original_error.context_exhausted
            raise original_error
        try:
            return self.normalize_plan(_loads_json_object(repaired_content), slide_count)
        except PptPlanParseError as repair_error:
            original_error.repair_error = repair_error.to_detail()
            original_error.context_exhausted = (
                context_exhausted
                or original_error.context_exhausted
                or (not _clean(text_base_url) and _looks_like_truncated_json(repaired_content, repair_error.attempts))
            )
            raise original_error from repair_error

    def _chat_text(self, payload: dict[str, Any], *, stream_label: str, base_url: str = "", api_key: str = "") -> str:
        if _clean(base_url):
            if config.ppt_plan_stream_log_enabled:
                content = self._collect_stream_text(
                    _post_openai_stream_json(base_url, api_key, "/v1/chat/completions", {**payload, "stream": True}),
                    stream_label=stream_label,
                )
                if content:
                    return content
            result = _post_openai_json(base_url, api_key, "/v1/chat/completions", {**payload, "stream": False})
            return _chat_content(result)

        if config.ppt_plan_stream_log_enabled:
            stream_payload = dict(payload)
            stream_payload["stream"] = True
            result = self.chat_handler(stream_payload)
            if not isinstance(result, dict):
                return self._collect_stream_text(result, stream_label=stream_label)
            return _chat_content(result)

        result = self.chat_handler({**payload, "stream": False})
        if not isinstance(result, dict):
            return "".join(_chat_chunk_content(chunk) for chunk in result if isinstance(chunk, dict)).strip()
        return _chat_content(result)

    @staticmethod
    def _collect_stream_text(chunks: Iterable[dict[str, Any]], *, stream_label: str) -> str:
        parts: list[str] = []
        print(f"[ppt-plan-stream] begin {stream_label}", flush=True)
        try:
            for chunk in chunks:
                if not isinstance(chunk, dict):
                    continue
                content = _chat_chunk_content(chunk)
                if not content:
                    continue
                parts.append(content)
                print(content, end="", flush=True)
        finally:
            print(f"\n[ppt-plan-stream] end {stream_label}", flush=True)
        return "".join(parts).strip()

    def create_master_task(
        self,
        identity: dict[str, object],
        *,
        client_task_id: str,
        name: str = "",
        model: str = DEFAULT_IMAGE_MODEL,
        size: str | None = None,
        quality: str | None = None,
        account_type: str = DEFAULT_IMAGE_ACCOUNT_TYPE,
        concurrency: object = DEFAULT_IMAGE_CONCURRENCY,
        style_prompt: str = "",
        image_base_url: str = "",
        image_api_key: str = "",
        base_url: str = "",
    ) -> dict[str, Any]:
        task_id = _clean(client_task_id)
        if not task_id:
            raise ValueError("client_task_id is required")
        normalized_concurrency = normalize_image_concurrency(concurrency)
        owner = _owner_id(identity)
        key = _task_key(owner, task_id)
        clean_image_base_url = _clean(image_base_url)
        clean_style_prompt = _clean(style_prompt)
        now = _now_iso()
        with self._lock:
            task = self._tasks.get(key)
            if task is not None:
                return self._public_task(task)
            slides = [
                {
                    "slide_id": layout_type,
                    "title": MASTER_LAYOUT_LABELS[layout_type],
                    "layout_type": layout_type,
                    "original_prompt": self._master_prompt(layout_type),
                    "current_prompt": self._master_prompt(layout_type),
                    "final_prompt": self._master_final_prompt(
                        layout_type,
                        self._master_prompt(layout_type),
                        clean_style_prompt,
                    ),
                    "image_url": "",
                    "version": 1,
                    "status": SLIDE_STATUS_DRAFT,
                    "error": "",
                }
                for layout_type in MASTER_LAYOUT_ORDER
            ]
            slides[0]["reference_images"] = []
            task = {
                "id": task_id,
                "name": _clean(name) or "PPT 母版",
                "owner_id": owner,
                "task_type": TASK_TYPE_MASTER,
                "master_confirmed": False,
                "status": TASK_STATUS_DRAFT,
                "model": _clean(model, DEFAULT_IMAGE_MODEL),
                "account_type": normalize_image_account_type(account_type),
                "size": normalize_ppt_image_size(size),
                "quality": normalize_image_quality(quality),
                "concurrency": normalized_concurrency,
                "image_base_url": clean_image_base_url,
                "image_api_key": _clean(image_api_key) if clean_image_base_url else "",
                "slide_count": len(slides),
                "design_concept": "PPT 母版审阅",
                "global_style_prompt": _master_style_context(clean_style_prompt),
                "master_style_prompt": clean_style_prompt,
                "markdown": "",
                "markdown_file_name": "",
                "slides": slides,
                "created_at": now,
                "updated_at": now,
                "error": "",
            }
            self._tasks[key] = task
            self._save_locked()
        return self._public_task(task)

    def save_plan_task(
        self,
        identity: dict[str, object],
        *,
        client_task_id: str,
        plan: dict[str, Any],
        master_task_id: str = "",
        markdown: str = "",
        markdown_file_name: str = "",
        name: str = "",
    ) -> dict[str, Any]:
        task_id = _clean(client_task_id)
        if not task_id:
            raise ValueError("client_task_id is required")
        normalized_plan = self.normalize_plan(plan, plan.get("slide_count"))
        owner = _owner_id(identity)
        clean_master_task_id = _clean(master_task_id)
        clean_markdown = _clean(markdown)
        clean_markdown_file_name = _clean(markdown_file_name)
        first_slide_title = _clean(normalized_plan["slides"][0].get("title")) if normalized_plan["slides"] else ""
        clean_name = (_clean(name) or clean_markdown_file_name or first_slide_title or task_id)[:120]
        key = _task_key(owner, task_id)
        now = _now_iso()
        with self._lock:
            task = self._tasks.get(key)
            if task is not None:
                return self._public_task(task)
            master_slides = (
                self._master_snapshot_locked(owner, clean_master_task_id, require_confirmed=True)
                if clean_master_task_id
                else []
            )
            slides = [
                {
                    "slide_id": slide["slide_id"],
                    "title": slide["title"],
                    "layout_type": slide.get("layout_type") or MASTER_LAYOUT_SINGLE_COLUMN,
                    "chapter_no": _clean(slide.get("chapter_no")),
                    "chapter_title": _clean(slide.get("chapter_title")),
                    "reference_images": self._default_content_references(master_slides, slide.get("layout_type")),
                    "original_prompt": slide["slide_prompt"],
                    "current_prompt": slide["slide_prompt"],
                    "final_prompt": "",
                    "image_url": "",
                    "version": 1,
                    "status": SLIDE_STATUS_SUCCESS,
                    "error": "",
                }
                for slide in normalized_plan["slides"]
            ]
            task = {
                "id": task_id,
                "name": clean_name,
                "owner_id": owner,
                "task_type": TASK_TYPE_PLAN,
                "master_confirmed": False,
                "status": TASK_STATUS_SUCCESS,
                "model": DEFAULT_IMAGE_MODEL,
                "account_type": DEFAULT_IMAGE_ACCOUNT_TYPE,
                "size": DEFAULT_PPT_IMAGE_SIZE,
                "quality": DEFAULT_IMAGE_QUALITY,
                "concurrency": DEFAULT_IMAGE_CONCURRENCY,
                "image_base_url": "",
                "image_api_key": "",
                "slide_count": normalized_plan["slide_count"],
                "design_concept": normalized_plan["design_concept"],
                "global_style_prompt": normalized_plan["global_style_prompt"],
                "chapters": normalized_plan.get("chapters", []),
                "markdown": clean_markdown,
                "markdown_file_name": clean_markdown_file_name,
                "slides": slides,
                "created_at": now,
                "updated_at": now,
                "error": "",
            }
            if clean_master_task_id:
                task["master_task_id"] = clean_master_task_id
                task["master_slides"] = master_slides
            self._tasks[key] = task
            self._save_locked()
            return self._public_task(task)

    def update_plan_task(
        self,
        identity: dict[str, object],
        task_id: str,
        plan: dict[str, Any],
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        normalized_plan = self.normalize_plan(plan, plan.get("slide_count"))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) != TASK_TYPE_PLAN:
                raise ValueError("只能更新方案任务")
            old_slides = {
                _clean(slide.get("slide_id")): slide
                for slide in self._slides_locked(task)
            }
            next_slides: list[dict[str, Any]] = []
            for slide in normalized_plan["slides"]:
                slide_id = _clean(slide.get("slide_id"))
                old_slide = old_slides.get(slide_id) or {}
                current_prompt = _clean(slide.get("slide_prompt"))
                layout_type = _normalize_master_layout(slide.get("layout_type"))
                title = _clean(slide.get("title"))
                version = int(old_slide.get("version") or 1)
                if (
                    _clean(old_slide.get("title")) != title
                    or _normalize_master_layout(old_slide.get("layout_type")) != layout_type
                    or _clean(old_slide.get("current_prompt")) != current_prompt
                ):
                    version += 1
                next_slide = {
                    "slide_id": slide_id,
                    "title": title,
                    "layout_type": layout_type,
                    "chapter_no": _clean(slide.get("chapter_no")),
                    "chapter_title": _clean(slide.get("chapter_title")),
                    "reference_images": _normalize_reference_images(old_slide.get("reference_images")),
                    "original_prompt": _clean(old_slide.get("original_prompt"), current_prompt),
                    "current_prompt": current_prompt,
                    "final_prompt": "",
                    "image_url": "",
                    "version": version,
                    "status": SLIDE_STATUS_SUCCESS,
                    "error": "",
                }
                next_slides.append(next_slide)
            task["slide_count"] = normalized_plan["slide_count"]
            task["design_concept"] = normalized_plan["design_concept"]
            task["global_style_prompt"] = normalized_plan["global_style_prompt"]
            task["chapters"] = normalized_plan.get("chapters", [])
            task["slides"] = next_slides
            task["status"] = TASK_STATUS_SUCCESS
            task["error"] = ""
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def create_task(
        self,
        identity: dict[str, object],
        *,
        client_task_id: str,
        plan: dict[str, Any],
        master_task_id: str = "",
        markdown: str = "",
        markdown_file_name: str = "",
        name: str = "",
        model: str = DEFAULT_IMAGE_MODEL,
        size: str | None = None,
        quality: str | None = None,
        account_type: str = DEFAULT_IMAGE_ACCOUNT_TYPE,
        concurrency: object = DEFAULT_IMAGE_CONCURRENCY,
        image_base_url: str = "",
        image_api_key: str = "",
        base_url: str = "",
    ) -> dict[str, Any]:
        task_id = _clean(client_task_id)
        if not task_id:
            raise ValueError("client_task_id is required")
        normalized_concurrency = normalize_image_concurrency(concurrency)
        owner = _owner_id(identity)
        key = _task_key(owner, task_id)
        normalized_plan = self.normalize_plan(plan, plan.get("slide_count"))
        clean_image_base_url = _clean(image_base_url)
        clean_master_task_id = _clean(master_task_id)
        clean_markdown = _clean(markdown)
        clean_markdown_file_name = _clean(markdown_file_name)
        first_slide_title = _clean(normalized_plan["slides"][0].get("title")) if normalized_plan["slides"] else ""
        clean_name = _clean(name) or clean_markdown_file_name or first_slide_title or task_id
        now = _now_iso()
        with self._lock:
            task = self._tasks.get(key)
            if task is not None:
                return self._public_task(task)
            master_slides = (
                self._master_snapshot_locked(owner, clean_master_task_id, require_confirmed=True)
                if clean_master_task_id
                else []
            )
            slides = [
                {
                    "slide_id": slide["slide_id"],
                    "title": slide["title"],
                    "layout_type": slide.get("layout_type") or MASTER_LAYOUT_SINGLE_COLUMN,
                    "chapter_no": _clean(slide.get("chapter_no")),
                    "chapter_title": _clean(slide.get("chapter_title")),
                    "reference_images": self._default_content_references(master_slides, slide.get("layout_type")),
                    "original_prompt": slide["slide_prompt"],
                    "current_prompt": slide["slide_prompt"],
                    "final_prompt": "",
                    "image_url": "",
                    "version": 1,
                    "status": SLIDE_STATUS_QUEUED,
                    "error": "",
                }
                for slide in normalized_plan["slides"]
            ]
            task = {
                "id": task_id,
                "name": clean_name,
                "owner_id": owner,
                "task_type": TASK_TYPE_CONTENT,
                "status": TASK_STATUS_QUEUED,
                "model": _clean(model, DEFAULT_IMAGE_MODEL),
                "account_type": normalize_image_account_type(account_type),
                "size": normalize_ppt_image_size(size),
                "quality": normalize_image_quality(quality),
                "concurrency": normalized_concurrency,
                "image_base_url": clean_image_base_url,
                "image_api_key": _clean(image_api_key) if clean_image_base_url else "",
                "slide_count": normalized_plan["slide_count"],
                "design_concept": normalized_plan["design_concept"],
                "global_style_prompt": normalized_plan["global_style_prompt"],
                "chapters": normalized_plan.get("chapters", []),
                "markdown": clean_markdown,
                "markdown_file_name": clean_markdown_file_name,
                "slides": slides,
                "created_at": now,
                "updated_at": now,
                "error": "",
            }
            if clean_master_task_id:
                task["master_task_id"] = clean_master_task_id
                task["master_slides"] = master_slides
            for slide in slides:
                slide["final_prompt"] = self._final_prompt(
                    normalized_plan["global_style_prompt"],
                    _clean(slide.get("current_prompt")),
                    layout_type=slide.get("layout_type") or MASTER_LAYOUT_SINGLE_COLUMN,
                    uses_master=bool(master_slides),
                    control_prompt=_slide_generation_control_prompt(task, slide),
                )
            self._tasks[key] = task
            self._save_locked()
        thread = threading.Thread(
            target=self._run_task,
            args=(key, base_url),
            name=f"ppt-task-{task_id[:16]}",
            daemon=True,
        )
        thread.start()
        return self._public_task(task)

    def list_tasks(self, identity: dict[str, object], task_ids: list[str]) -> dict[str, Any]:
        owner = _owner_id(identity)
        requested_ids = [_clean(task_id) for task_id in task_ids if _clean(task_id)]
        with self._lock:
            items = []
            missing_ids = []
            for task_id in requested_ids:
                task = self._tasks.get(_task_key(owner, task_id))
                if task is None:
                    missing_ids.append(task_id)
                else:
                    items.append(self._public_task(task))
            if not requested_ids:
                items = [
                    self._public_task(task)
                    for task in self._tasks.values()
                    if task.get("owner_id") == owner
                ]
                items.sort(key=lambda item: str(item.get("updated_at") or ""), reverse=True)
                missing_ids = []
            return {"items": items, "missing_ids": missing_ids}

    def require_master_ready(
        self,
        identity: dict[str, object],
        master_task_id: str,
        *,
        require_confirmed: bool = True,
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        clean_master_task_id = _clean(master_task_id)
        with self._lock:
            task = self._tasks.get(_task_key(owner, clean_master_task_id))
            self._validate_master_task_locked(task, clean_master_task_id, require_confirmed=require_confirmed)
            return self._public_task(task)

    def confirm_master_task(self, identity: dict[str, object], task_id: str) -> dict[str, Any]:
        owner = _owner_id(identity)
        clean_task_id = _clean(task_id)
        key = _task_key(owner, clean_task_id)
        with self._lock:
            task = self._tasks.get(key)
            self._validate_master_task_locked(task, clean_task_id, require_confirmed=False)
            task["master_confirmed"] = True
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def rename_task(self, identity: dict[str, object], task_id: str, name: str) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_name = _clean(name)
        if not clean_name:
            raise ValueError("name is required")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            task["name"] = clean_name[:120]
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def test_provider(self, *, kind: str, model: str = "", base_url: str = "", api_key: str = "") -> dict[str, Any]:
        provider_kind = _clean(kind).lower()
        if provider_kind not in {"text", "image"}:
            raise ValueError("provider kind must be text or image")
        started = time.time()
        target_model = _clean(model, DEFAULT_TEXT_MODEL if provider_kind == "text" else DEFAULT_IMAGE_MODEL)
        clean_base_url = _clean(base_url)
        mode = "external" if clean_base_url else "current_project"
        if clean_base_url:
            status, result = _get_openai_json(clean_base_url, api_key, "/v1/models", timeout=10.0)
        else:
            status = 200
            result = openai_v1_models.list_models()
        ids = _model_ids(result)
        model_found = target_model in ids if target_model else False
        if ids and target_model and not model_found:
            message = f"服务可访问，但 /v1/models 未返回模型 {target_model}"
        else:
            message = "服务可访问"
        return {
            "ok": True,
            "kind": provider_kind,
            "mode": mode,
            "status": status,
            "latency_ms": int((time.time() - started) * 1000),
            "model": target_model,
            "model_found": model_found,
            "model_count": len(ids),
            "message": message,
        }

    def delete_task(self, identity: dict[str, object], task_id: str) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        with self._lock:
            task = self._tasks.pop(key, None)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            pptx_path = _clean(task.get("pptx_path"))
            self._save_locked()
        if pptx_path:
            try:
                Path(pptx_path).unlink(missing_ok=True)
            except OSError:
                pass
        return {"ok": True}

    def stop_task(self, identity: dict[str, object], task_id: str) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_PLAN:
                raise ValueError("方案任务不需要停止")
            has_active_work = task.get("status") in UNFINISHED_TASK_STATUSES or any(
                slide.get("status") in UNFINISHED_SLIDE_STATUSES for slide in self._slides_locked(task)
            )
            if not has_active_work:
                return self._public_task(task)
            task["cancel_requested"] = True
            task["status"] = TASK_STATUS_STOPPED
            task["error"] = "任务已停止"
            task["updated_at"] = _now_iso()
            task["finished_at"] = task["updated_at"]
            self._invalidate_package_locked(task)
            for slide in self._slides_locked(task):
                if slide.get("status") in UNFINISHED_SLIDE_STATUSES:
                    slide["status"] = SLIDE_STATUS_STOPPED
                    slide["error"] = "已停止"
                    slide["finished_at"] = task["updated_at"]
            self._save_locked()
            return self._public_task(task)

    def resume_task(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        concurrency: object | None = None,
        model: str = "",
        account_type: str | None = None,
        size: str | None = None,
        quality: str | None = None,
        image_base_url: str | None = None,
        image_api_key: str | None = None,
        base_url: str = "",
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            if task.get("status") in UNFINISHED_TASK_STATUSES:
                return self._public_task(task)
            is_master_task = _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER
            task.pop("cancel_requested", None)
            task.pop("finished_at", None)
            if concurrency is not None:
                task["concurrency"] = normalize_image_concurrency(concurrency)
            if _clean(model):
                task["model"] = _clean(model, DEFAULT_IMAGE_MODEL)
            if account_type is not None:
                task["account_type"] = normalize_image_account_type(account_type)
            if size is not None:
                task["size"] = normalize_ppt_image_size(size)
            if quality is not None:
                task["quality"] = normalize_image_quality(quality)
            next_image_base_url = _clean(task.get("image_base_url"))
            if image_base_url is not None:
                next_image_base_url = _clean(image_base_url)
                task["image_base_url"] = next_image_base_url
            if image_api_key is not None:
                task["image_api_key"] = _clean(image_api_key) if next_image_base_url else ""
            if not next_image_base_url:
                task["image_api_key"] = ""
            queued: list[tuple[str, int]] = []
            for slide in self._slides_locked(task):
                if slide.get("status") == SLIDE_STATUS_SUCCESS and _clean(slide.get("image_url")):
                    continue
                prompt_source = _strip_generated_chapter_instructions(_clean(slide.get("current_prompt") or slide.get("original_prompt")))
                if not prompt_source and not _clean(slide.get("final_prompt")):
                    slide["status"] = SLIDE_STATUS_ERROR
                    slide["error"] = "请输入图片描述后生成图片"
                    continue
                version = int(slide.get("version") or 1) + 1
                slide["version"] = version
                slide["status"] = SLIDE_STATUS_QUEUED
                slide["error"] = ""
                slide.pop("finished_at", None)
                layout_type = _normalize_master_layout(slide.get("layout_type"))
                slide["layout_type"] = layout_type
                if is_master_task:
                    if not _clean(slide.get("final_prompt")):
                        slide["final_prompt"] = self._master_final_prompt(
                            layout_type,
                            prompt_source,
                            task.get("master_style_prompt"),
                        )
                else:
                    slide["final_prompt"] = self._final_prompt(
                        _clean(task.get("global_style_prompt"), GLOBAL_STYLE_PROMPT),
                        prompt_source,
                        layout_type=layout_type,
                        uses_master=bool(task.get("master_slides") or _clean(task.get("master_task_id"))),
                        control_prompt=_slide_generation_control_prompt(task, slide),
                    )
                queued.append((_clean(slide.get("slide_id")), version))
            if not queued:
                self._finalize_task_status(key)
                return self._public_task(task)
            task["status"] = TASK_STATUS_RUNNING
            task["error"] = ""
            task.pop("pptx_path", None)
            task.pop("packaged_at", None)
            task["started_at"] = _now_iso()
            task["updated_at"] = _now_iso()
            self._save_locked()
        thread = threading.Thread(
            target=self._run_master_task if is_master_task else self._run_task,
            args=(key, base_url),
            name=f"ppt-resume-{_clean(task_id)[:16]}",
            daemon=True,
        )
        thread.start()
        with self._lock:
            return self._public_task(self._tasks[key])

    def regenerate_slide(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        prompt: str,
        base_url: str = "",
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_prompt = _strip_generated_chapter_instructions(prompt)
        if not clean_prompt:
            raise ValueError("prompt is required")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            layout_type = _normalize_master_layout(slide.get("layout_type"))
            is_master_task = _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER
            version = int(slide.get("version") or 1) + 1
            slide["version"] = version
            slide["layout_type"] = layout_type
            slide["current_prompt"] = clean_prompt
            if is_master_task:
                slide["final_prompt"] = self._master_final_prompt(
                    layout_type,
                    clean_prompt,
                    task.get("master_style_prompt"),
                )
            else:
                slide["final_prompt"] = self._final_prompt(
                    _clean(task.get("global_style_prompt"), GLOBAL_STYLE_PROMPT),
                    clean_prompt,
                    layout_type=layout_type,
                    uses_master=bool(task.get("master_slides") or _clean(task.get("master_task_id"))),
                    control_prompt=_slide_generation_control_prompt(task, slide),
                )
            slide["status"] = SLIDE_STATUS_QUEUED
            slide["error"] = ""
            slide.pop("finished_at", None)
            task.pop("cancel_requested", None)
            task.pop("finished_at", None)
            task["status"] = TASK_STATUS_RUNNING
            task["error"] = ""
            task["started_at"] = _now_iso()
            if is_master_task:
                task["master_confirmed"] = False
            task.pop("pptx_path", None)
            task.pop("packaged_at", None)
            task["updated_at"] = _now_iso()
            self._save_locked()
        thread = threading.Thread(
            target=self._run_slide,
            args=(key, clean_slide_id, version, base_url),
            name=f"ppt-slide-{_clean(task_id)[:16]}-{clean_slide_id}",
            daemon=True,
        )
        thread.start()
        with self._lock:
            return self._public_task(self._tasks[key])

    def update_slide_prompt(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        prompt: str,
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_prompt = _strip_generated_chapter_instructions(prompt)
        if not clean_prompt:
            raise ValueError("prompt is required")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            if task.get("status") in UNFINISHED_TASK_STATUSES:
                raise ValueError("图片生成过程中不允许修改提示词")
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            layout_type = _normalize_master_layout(slide.get("layout_type"))
            is_master_task = _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER
            if _clean(slide.get("current_prompt")) != clean_prompt:
                slide["version"] = int(slide.get("version") or 1) + 1
            slide["layout_type"] = layout_type
            slide["current_prompt"] = clean_prompt
            if is_master_task:
                slide["final_prompt"] = self._master_final_prompt(
                    layout_type,
                    clean_prompt,
                    task.get("master_style_prompt"),
                )
            else:
                slide["final_prompt"] = self._final_prompt(
                    _clean(task.get("global_style_prompt"), GLOBAL_STYLE_PROMPT),
                    clean_prompt,
                    layout_type=layout_type,
                    uses_master=bool(task.get("master_slides") or _clean(task.get("master_task_id"))),
                    control_prompt=_slide_generation_control_prompt(task, slide),
                )
            if is_master_task:
                task["master_confirmed"] = False
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def upload_slide_image(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        image_url: str,
        base_url: str = "",
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_image_url = _clean(image_url)
        if not clean_image_url:
            raise ValueError("image_url is required")
        image_metadata = _image_metadata_from_bytes(self.image_fetcher(clean_image_url, base_url))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            slide["image_url"] = clean_image_url
            slide["status"] = SLIDE_STATUS_SUCCESS
            slide["error"] = ""
            slide.update(image_metadata)
            slide["finished_at"] = _now_iso()
            slide["version"] = int(slide.get("version") or 1) + 1
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER:
                task["master_confirmed"] = False
            self._invalidate_package_locked(task)
            self._finalize_task_status_locked(task)
            self._save_locked()
            return self._public_task(task)

    def add_slide_reference(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        image_url: str,
        title: str = "",
        reference_id: str = "",
        base_url: str = "",
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_image_url = _clean(image_url)
        if not clean_image_url:
            raise ValueError("image_url is required")
        self.image_fetcher(clean_image_url, base_url)
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            if task.get("status") in UNFINISHED_TASK_STATUSES or any(
                item.get("status") in UNFINISHED_SLIDE_STATUSES for item in self._slides_locked(task)
            ):
                raise ValueError("图片生成过程中不允许修改参考图")
            if "reference_images" in slide:
                references = _normalize_reference_images(slide.get("reference_images"))
            else:
                try:
                    references = self._slide_reference_images_locked(task, slide)
                except Exception:
                    references = []
            clean_reference_id = _clean(reference_id) or f"ref-{uuid.uuid4().hex[:12]}"
            references.append({
                "id": clean_reference_id,
                "title": _clean(title, "用户参考图"),
                "layout_type": _normalize_master_layout(slide.get("layout_type") or slide.get("slide_id")),
                "image_url": clean_image_url,
            })
            slide["reference_images"] = references
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER:
                task["master_confirmed"] = False
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def delete_slide_reference(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        reference_id: str,
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_reference_id = _clean(reference_id)
        if not clean_reference_id:
            raise ValueError("reference_id is required")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            if task.get("status") in UNFINISHED_TASK_STATUSES or any(
                item.get("status") in UNFINISHED_SLIDE_STATUSES for item in self._slides_locked(task)
            ):
                raise ValueError("图片生成过程中不允许删除参考图")
            if "reference_images" in slide:
                references = _normalize_reference_images(slide.get("reference_images"))
            else:
                references = self._slide_reference_images_locked(task, slide)
            next_references = [item for item in references if _clean(item.get("id")) != clean_reference_id]
            if len(next_references) == len(references):
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}/references/{reference_id}")
            slide["reference_images"] = next_references
            task["updated_at"] = _now_iso()
            self._save_locked()
            return self._public_task(task)

    def edit_slide_image(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        prompt: str,
        base_url: str = "",
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_prompt = _clean(prompt)
        if not clean_prompt:
            raise ValueError("prompt is required")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            reference_image_url = _clean(slide.get("image_url"))
            if not reference_image_url:
                raise ValueError("请先上传或生成本页图片")
            version = int(slide.get("version") or 1) + 1
            slide["version"] = version
            slide["status"] = SLIDE_STATUS_QUEUED
            slide["error"] = ""
            slide.pop("finished_at", None)
            task.pop("cancel_requested", None)
            task.pop("finished_at", None)
            task["status"] = TASK_STATUS_RUNNING
            task["error"] = ""
            task["started_at"] = _now_iso()
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER:
                task["master_confirmed"] = False
            self._invalidate_package_locked(task)
            task["updated_at"] = _now_iso()
            self._save_locked()
        thread = threading.Thread(
            target=self._run_slide_edit,
            args=(key, clean_slide_id, version, reference_image_url, clean_prompt, base_url),
            name=f"ppt-edit-{_clean(task_id)[:16]}-{clean_slide_id}",
            daemon=True,
        )
        thread.start()
        with self._lock:
            return self._public_task(self._tasks[key])

    def insert_blank_slide(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
        position: str,
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        clean_position = _clean(position).lower()
        if clean_position not in {"before", "after"}:
            raise ValueError("position must be before or after")
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slides = self._slides_locked(task)
            if len(slides) >= MAX_SLIDE_COUNT:
                raise ValueError(f"slide_count must be an integer between {MIN_SLIDE_COUNT} and {MAX_SLIDE_COUNT}")
            target_index = next((index for index, slide in enumerate(slides) if _clean(slide.get("slide_id")) == clean_slide_id), -1)
            if target_index < 0:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            insert_index = target_index if clean_position == "before" else target_index + 1
            new_slide_id = f"blank-{uuid.uuid4().hex[:10]}"
            slides.insert(
                insert_index,
                {
                    "slide_id": new_slide_id,
                    "title": "空白页",
                    "layout_type": MASTER_LAYOUT_SINGLE_COLUMN,
                    "reference_images": [],
                    "original_prompt": "",
                    "current_prompt": "",
                    "final_prompt": "",
                    "image_url": "",
                    "version": 1,
                    "status": SLIDE_STATUS_ERROR,
                    "error": "请输入图片描述后生成图片",
                },
            )
            task["slide_count"] = len(slides)
            self._invalidate_package_locked(task)
            self._finalize_task_status_locked(task)
            self._save_locked()
            return self._public_task(task)

    def delete_slide(
        self,
        identity: dict[str, object],
        *,
        task_id: str,
        slide_id: str,
    ) -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        clean_slide_id = _clean(slide_id)
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slides = self._slides_locked(task)
            if len(slides) <= 1:
                raise ValueError("至少需要保留 1 页 PPT")
            target_index = next((index for index, slide in enumerate(slides) if _clean(slide.get("slide_id")) == clean_slide_id), -1)
            if target_index < 0:
                raise PptTaskNotFoundError(f"{task_id}/{slide_id}")
            slides.pop(target_index)
            task["slide_count"] = len(slides)
            self._invalidate_package_locked(task)
            self._finalize_task_status_locked(task)
            self._save_locked()
            return self._public_task(task)

    def package_task(self, identity: dict[str, object], *, task_id: str, base_url: str = "") -> dict[str, Any]:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            slides = self._slides_locked(task)
            if any(slide.get("status") in UNFINISHED_SLIDE_STATUSES for slide in slides):
                raise ValueError("仍有页面正在生成，暂时不能打包 PPTX")
            if not slides or any(not _clean(slide.get("image_url")) for slide in slides):
                raise ValueError("所有页面都有成功图片后才能打包 PPTX")
            task["status"] = TASK_STATUS_PACKAGING
            task["error"] = ""
            task["updated_at"] = _now_iso()
            self._save_locked()
            snapshot = [(str(slide.get("title") or ""), str(slide.get("image_url") or "")) for slide in slides]
        try:
            image_slides = [(title, self.image_fetcher(image_url, base_url)) for title, image_url in snapshot]
            payload = build_image_pptx(image_slides)
            package_path = self._package_path(owner, _clean(task_id))
            package_path.parent.mkdir(parents=True, exist_ok=True)
            package_path.write_bytes(payload)
            with self._lock:
                task = self._tasks.get(key)
                if task is None:
                    raise PptTaskNotFoundError(task_id)
                task["status"] = TASK_STATUS_PACKAGED
                task["pptx_path"] = str(package_path)
                task["pptx_build_version"] = PPTX_BUILD_VERSION
                task["packaged_at"] = _now_iso()
                task["updated_at"] = task["packaged_at"]
                task["error"] = ""
                self._save_locked()
                return self._public_task(task)
        except Exception as exc:
            with self._lock:
                task = self._tasks.get(key)
                if task is not None:
                    task["status"] = TASK_STATUS_ERROR
                    task["error"] = str(exc) or "打包 PPTX 失败"
                    task["updated_at"] = _now_iso()
                    self._save_locked()
            raise

    def download_path(self, identity: dict[str, object], task_id: str) -> Path:
        owner = _owner_id(identity)
        key = _task_key(owner, _clean(task_id))
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            pptx_path = _clean(task.get("pptx_path"))
            pptx_build_version = _pptx_build_version(task.get("pptx_build_version"))
        if not pptx_path:
            raise ValueError("PPTX 尚未打包")
        if pptx_build_version != PPTX_BUILD_VERSION:
            raise ValueError("PPTX 文件格式已更新，请重新打包后下载")
        path = Path(pptx_path)
        if not path.is_file():
            raise ValueError("PPTX 文件不存在，请重新打包")
        return path

    def download_slide_image(self, identity: dict[str, object], task_id: str, slide_id: str, base_url: str = "") -> tuple[bytes, str, str]:
        owner = _owner_id(identity)
        clean_task_id = _clean(task_id)
        clean_slide_id = _clean(slide_id)
        key = _task_key(owner, clean_task_id)
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                raise PptTaskNotFoundError(task_id)
            candidates = self._slides_locked(task)
            slide = self._find_slide_locked(task, clean_slide_id)
            if slide is None:
                master_slides = task.get("master_slides")
                if isinstance(master_slides, list):
                    candidates = [item for item in master_slides if isinstance(item, dict)]
                    slide = next((item for item in candidates if _clean(item.get("slide_id")) == clean_slide_id), None)
            if slide is None:
                raise ValueError("PPT 页面不存在")
            image_url = _clean(slide.get("image_url"))
            title = _clean(slide.get("title"), clean_slide_id)
            display_index = next((index for index, item in enumerate(candidates, start=1) if _clean(item.get("slide_id")) == clean_slide_id), 1)
        if not image_url:
            raise ValueError("PPT 页面没有可下载的图片")
        data = self.image_fetcher(image_url, base_url)
        mime_type = _mime_from_image_bytes(data)
        filename = f"ppt-{display_index}-{_safe_file_stem(title, clean_slide_id or 'slide')}.{_extension_from_mime(mime_type)}"
        return data, mime_type, filename

    def _plan_messages(self, markdown: str, slide_count: int | None) -> list[dict[str, str]]:
        schema = {
            "slide_count": slide_count or 8,
            "design_concept": "整体设计思路",
            "global_style_prompt": GLOBAL_STYLE_PROMPT,
            "layout_options": MASTER_LAYOUT_LABELS,
            "chapters": [
                {
                    "chapter_no": "01",
                    "title": "章节标题",
                    "section_slide_id": "3",
                }
            ],
            "slides": [
                {
                    "slide_id": "1",
                    "title": "页面标题",
                    "layout_type": MASTER_LAYOUT_SINGLE_COLUMN,
                    "chapter_no": "01",
                    "chapter_title": "章节标题",
                    "slide_prompt": "描述该页核心信息、视觉层级、图表/构图重点和需要呈现的文字",
                }
            ],
        }
        return [
            {
                "role": "system",
                "content": (
                    "你是资深 PPT 信息架构和视觉设计顾问。根据用户 Markdown 生成可用于 AI 生图的逐页设计方案。"
                    "必须只输出合法 JSON，不要输出 Markdown 代码块或解释。"
                    f"每页必须从 {', '.join(MASTER_LAYOUT_ORDER)} 中选择一个 layout_type。"
                    f"{NO_PAGE_NUMBER_PROMPT}"
                ),
            },
            {
                "role": "user",
                "content": (
                    (
                        f"请把下面的 Markdown 精确规划为 {slide_count} 页 PPT。"
                        if slide_count is not None
                        else f"请根据下面 Markdown 的信息密度自行决定 PPT 页数，建议 {MIN_SLIDE_COUNT}-{MAX_SLIDE_COUNT} 页，并让 slide_count 与 slides 数量完全一致。"
                    )
                    + "每页 slide_prompt 需要可直接交给图片模型基于母版填充一张完整页面，包含标题、核心信息、构图、图表和视觉重点。"
                    "必须额外输出 chapters 数组作为唯一章节清单；目录页和所有 section_break 页都必须使用 chapters 中完全一致的章节编号和章节标题。"
                    "每页必须输出 layout_type：cover 表示封面页；agenda 表示目录页；section_break 表示章节过渡页；"
                    "single_column 表示单栏内容页；two_column 表示双栏图文页；bento_card 表示卡片布局页；"
                    "dashboard 表示数据图表页；thank_you 表示结束页。"
                    "封面页必须提供整套 PPT 的主题；目录页必须提供 N 个章节标题；"
                    "章节过渡页必须说明当前是第几个章节以及章节标题；"
                    "结束页必须简洁，不允许包含联系电话、电子信箱、二维码、网址或任何联系方式。"
                    f"{PPT_STRUCTURE_PROMPT}"
                    f"{NO_PAGE_NUMBER_PROMPT}"
                    "所有页面文字都必须遵守字体规范：中文微软雅黑，英文 Times New Roman。"
                    f"{TITLE_STYLE_PROMPT}"
                    f"全局风格固定为：{GLOBAL_STYLE_PROMPT}\n"
                    "JSON 结构示例：\n"
                    f"{json.dumps(schema, ensure_ascii=False, indent=2)}\n\n"
                    "Markdown 原文：\n"
                    f"{markdown}"
                ),
            },
        ]

    @staticmethod
    def _final_prompt(
        global_style_prompt: str,
        slide_prompt: str,
        *,
        layout_type: str = "",
        uses_master: bool = False,
        control_prompt: str = "",
    ) -> str:
        parts = [
            _clean(global_style_prompt, GLOBAL_STYLE_PROMPT),
            FONT_STYLE_PROMPT,
            TITLE_STYLE_PROMPT,
        ]
        layout = _normalize_master_layout(layout_type) if layout_type else ""
        if uses_master:
            parts.append(
                "母版引用规范：输入图片是已确认的 PPT 母版，必须沿用其背景、版式、视觉层级、"
                "安全边距、装饰元素和空间关系；只替换或填充本页需要的文字、图表和内容，不要重新设计模板。"
            )
        if layout:
            parts.append(f"母版类型：{layout}（{_master_layout_label(layout)}）。{_master_layout_instruction(layout)}")
        parts.extend([
            _clean(control_prompt),
            _strip_generated_chapter_instructions(slide_prompt),
            FINAL_IMAGE_PROMPT_SUFFIX,
        ])
        return "\n\n".join(part for part in parts if _clean(part))

    @staticmethod
    def _master_prompt(layout_type: str) -> str:
        layout = _normalize_master_layout(layout_type)
        prompts = {
            MASTER_LAYOUT_COVER: (
                "封面页（Cover）：画面不要太满，保留大留白；标题区域要明显，方便后期改字。"
                "只设置标题、副标题、日期/单位等可替换占位区域，不要堆满装饰元素。"
            ),
            MASTER_LAYOUT_AGENDA: (
                "目录页（Agenda）：强调层级感和对齐，留出编号/章节位置；"
                "目录列表清晰，章节序号、标题和短说明有稳定对齐关系。"
            ),
            MASTER_LAYOUT_SECTION_BREAK: (
                "章节过渡页（Section Break）：视觉冲击强，但元素少；"
                "大标题居中或突出，章节编号区域明确，只保留少量辅助元素；同类章节页编号和标题字号保持一致。"
            ),
            MASTER_LAYOUT_SINGLE_COLUMN: (
                "单栏内容页（Single Column）：内容区不要复杂，保持统一边距和留白；"
                "预留标题、正文要点、图表或说明的单一主内容区域；一级标题位置和字号固定。"
            ),
            MASTER_LAYOUT_TWO_COLUMN: (
                "双栏图文页（Two-column）：左右比例平衡，图片区和文字区分明确；"
                "两栏有清晰边界、统一边距和可替换图片区；一级标题位置和字号固定。"
            ),
            MASTER_LAYOUT_BENTO_CARD: (
                "卡片布局页（Bento / Card）：卡片大小有节奏变化，注意统一圆角、阴影、间距；"
                "保留多块可替换内容卡片，避免卡片过密；一级标题位置和字号固定。"
            ),
            MASTER_LAYOUT_DASHBOARD: (
                "数据图表页（Dashboard）：图表区域留空，避免 AI 画死数据；"
                "KPI 数字区域要大，保留图表容器、指标卡和说明位置；一级标题位置和字号固定。"
            ),
            MASTER_LAYOUT_THANK_YOU: (
                "结束页（Thank You / Q&A）：尽量简洁，只保留感谢语或 Q&A 标题；"
                "不要设置联系电话、邮箱、二维码、网址或联系信息占位，不要加入复杂图表。"
            ),
        }
        return prompts.get(layout, prompts[MASTER_LAYOUT_SINGLE_COLUMN])

    @staticmethod
    def _master_final_prompt(layout_type: str, visible_prompt: str, style_prompt: object = "") -> str:
        prompt = _clean(visible_prompt) or PptTaskService._master_prompt(layout_type)
        return PptTaskService._final_prompt(
            _master_style_context(style_prompt),
            f"{MASTER_TEMPLATE_BASE_PROMPT}\n{prompt}",
            layout_type=layout_type,
        )

    def _run_master_task(self, key: str, base_url: str) -> None:
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                return
            task["status"] = TASK_STATUS_RUNNING
            now = _now_iso()
            task["started_at"] = task.get("started_at") or now
            task["updated_at"] = now
            base_slide = self._find_slide_locked(task, MASTER_BASE_LAYOUT)
            if base_slide is None:
                task["status"] = TASK_STATUS_ERROR
                task["error"] = "母版任务缺少封面页"
                self._save_locked()
                return
            base_version = int(base_slide.get("version") or 1)
            self._save_locked()

        self._run_slide(key, MASTER_BASE_LAYOUT, base_version, base_url, False)

        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                self._save_locked()
                return
            base_slide = self._find_slide_locked(task, MASTER_BASE_LAYOUT)
            base_url_ref = _clean(base_slide.get("image_url") if base_slide else "")
            base_ok = bool(base_slide and base_slide.get("status") == SLIDE_STATUS_SUCCESS and base_url_ref)
            if not base_ok:
                for slide in self._slides_locked(task):
                    if _clean(slide.get("slide_id")) == MASTER_BASE_LAYOUT:
                        continue
                    if slide.get("status") in UNFINISHED_SLIDE_STATUSES:
                        slide["status"] = SLIDE_STATUS_ERROR
                        slide["error"] = "封面页母版生成失败，无法继续生成其他母版"
                self._finalize_task_status_locked(task)
                self._save_locked()
                return
            base_reference = _reference_image_from_slide(base_slide)
            for slide in self._slides_locked(task):
                if _clean(slide.get("slide_id")) == MASTER_BASE_LAYOUT:
                    slide["reference_images"] = []
                elif "reference_images" not in slide and slide.get("status") == SLIDE_STATUS_QUEUED:
                    slide["reference_images"] = [base_reference]
            slide_versions = [
                (_clean(slide.get("slide_id")), int(slide.get("version") or 1))
                for slide in self._slides_locked(task)
                if _clean(slide.get("slide_id")) != MASTER_BASE_LAYOUT and slide.get("status") == SLIDE_STATUS_QUEUED
            ]
            concurrency = normalize_image_concurrency(task.get("concurrency"))
        if slide_versions:
            max_workers = min(concurrency, len(slide_versions))
            with ThreadPoolExecutor(max_workers=max_workers, thread_name_prefix="ppt-master-slide") as executor:
                futures = [
                    executor.submit(self._run_slide, key, slide_id, version, base_url, False, base_url_ref)
                    for slide_id, version in slide_versions
                ]
                for future in as_completed(futures):
                    future.result()
        self._finalize_task_status(key)

    def _run_task(self, key: str, base_url: str) -> None:
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                return
            task["status"] = TASK_STATUS_RUNNING
            now = _now_iso()
            task["started_at"] = task.get("started_at") or now
            task["updated_at"] = now
            self._save_locked()
            slide_versions = [
                (_clean(slide.get("slide_id")), int(slide.get("version") or 1))
                for slide in self._slides_locked(task)
                if slide.get("status") == SLIDE_STATUS_QUEUED
            ]
            concurrency = normalize_image_concurrency(task.get("concurrency"))
        if slide_versions:
            max_workers = min(concurrency, len(slide_versions))
            with ThreadPoolExecutor(max_workers=max_workers, thread_name_prefix="ppt-slide") as executor:
                futures = [
                    executor.submit(self._run_slide, key, slide_id, version, base_url, False)
                    for slide_id, version in slide_versions
                ]
                for future in as_completed(futures):
                    future.result()
        self._finalize_task_status(key)

    def _run_slide(
        self,
        key: str,
        slide_id: str,
        version: int,
        base_url: str,
        finalize: bool = True,
        reference_image_url: str = "",
    ) -> None:
        reference_error = ""
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                return
            slide = self._find_slide_locked(task, slide_id)
            if slide is None or int(slide.get("version") or 0) != version:
                return
            if slide.get("status") == SLIDE_STATUS_STOPPED:
                return
            slide["status"] = SLIDE_STATUS_RUNNING
            slide["error"] = ""
            slide["started_at"] = _now_iso()
            slide.pop("finished_at", None)
            task["status"] = TASK_STATUS_RUNNING
            task["updated_at"] = _now_iso()
            try:
                resolved_reference_images = (
                    [{"id": "reference", "title": "参考图", "layout_type": "", "image_url": _clean(reference_image_url)}]
                    if _clean(reference_image_url)
                    else self._slide_reference_images_locked(task, slide)
                )
            except Exception as exc:
                resolved_reference_images = []
                reference_error = str(exc) or "未找到母版参考图"
            payload = {
                "prompt": _clean(slide.get("final_prompt")),
                "model": _clean(task.get("model"), DEFAULT_IMAGE_MODEL),
                "account_type": normalize_image_account_type(task.get("account_type")),
                "n": 1,
                "size": normalize_ppt_image_size(task.get("size")),
                "quality": normalize_image_quality(task.get("quality")),
                "response_format": "url",
                "base_url": base_url,
            }
            image_base_url = _clean(task.get("image_base_url"))
            image_api_key = _clean(task.get("image_api_key"))
            self._save_locked()
        if reference_error:
            self._finish_slide(key, slide_id, version, status=SLIDE_STATUS_ERROR, error=reference_error)
            if finalize:
                self._finalize_task_status(key)
            return
        try:
            if resolved_reference_images:
                image_inputs = [
                    _image_input_from_url(_clean(reference.get("image_url")), base_url)
                    for reference in resolved_reference_images
                    if _clean(reference.get("image_url"))
                ]
                if not image_inputs:
                    raise ValueError("参考图地址为空")
                if image_base_url:
                    external_payload = {
                        payload_key: value
                        for payload_key, value in payload.items()
                        if payload_key not in {"base_url", "account_type"} and value is not None
                    }
                    external_payload["images"] = [{"image_url": _image_input_to_data_url(image_input)} for image_input in image_inputs]
                    result = _post_openai_json(image_base_url, image_api_key, "/v1/images/edits", external_payload)
                else:
                    result = self.image_edit_handler({**payload, "images": image_inputs})
            elif image_base_url:
                external_payload = {
                    payload_key: value
                    for payload_key, value in payload.items()
                    if payload_key not in {"base_url", "account_type"} and value is not None
                }
                result = _post_openai_json(image_base_url, image_api_key, "/v1/images/generations", external_payload)
            else:
                result = self.image_handler(payload)
            if not isinstance(result, dict):
                raise RuntimeError("image generation returned streaming result unexpectedly")
            image_url = _extract_image_url(result)
            image_metadata = _image_metadata_from_url(image_url, base_url)
            self._finish_slide(key, slide_id, version, status=SLIDE_STATUS_SUCCESS, image_url=image_url, error="", image_metadata=image_metadata)
        except Exception as exc:
            self._finish_slide(key, slide_id, version, status=SLIDE_STATUS_ERROR, error=str(exc) or "图片生成失败")
        if finalize:
            self._finalize_task_status(key)

    def _run_slide_edit(self, key: str, slide_id: str, version: int, reference_image_url: str, edit_prompt: str, base_url: str) -> None:
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                return
            slide = self._find_slide_locked(task, slide_id)
            if slide is None or int(slide.get("version") or 0) != version:
                return
            if slide.get("status") == SLIDE_STATUS_STOPPED:
                return
            slide["status"] = SLIDE_STATUS_RUNNING
            slide["error"] = ""
            slide["started_at"] = _now_iso()
            slide.pop("finished_at", None)
            task["status"] = TASK_STATUS_RUNNING
            task["updated_at"] = _now_iso()
            layout_type = _normalize_master_layout(slide.get("layout_type"))
            if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER:
                final_edit_prompt = self._master_final_prompt(layout_type, edit_prompt, task.get("master_style_prompt"))
            else:
                final_edit_prompt = self._final_prompt(
                    _clean(task.get("global_style_prompt"), GLOBAL_STYLE_PROMPT),
                    edit_prompt,
                    layout_type=layout_type,
                    uses_master=bool(task.get("master_slides") or _clean(task.get("master_task_id"))),
                    control_prompt=_slide_generation_control_prompt(task, slide),
                )
            payload = {
                "prompt": final_edit_prompt,
                "model": _clean(task.get("model"), DEFAULT_IMAGE_MODEL),
                "account_type": normalize_image_account_type(task.get("account_type")),
                "n": 1,
                "size": normalize_ppt_image_size(task.get("size")),
                "quality": normalize_image_quality(task.get("quality")),
                "response_format": "url",
                "base_url": base_url,
            }
            image_base_url = _clean(task.get("image_base_url"))
            image_api_key = _clean(task.get("image_api_key"))
            self._save_locked()
        try:
            image_input = _image_input_from_url(reference_image_url, base_url)
            if image_base_url:
                external_payload = {
                    payload_key: value
                    for payload_key, value in payload.items()
                    if payload_key not in {"base_url", "account_type"} and value is not None
                }
                external_payload["images"] = [{"image_url": _image_input_to_data_url(image_input)}]
                result = _post_openai_json(image_base_url, image_api_key, "/v1/images/edits", external_payload)
            else:
                result = self.image_edit_handler({**payload, "images": [image_input]})
            if not isinstance(result, dict):
                raise RuntimeError("image edit returned streaming result unexpectedly")
            image_url = _extract_image_url(result)
            image_metadata = _image_metadata_from_url(image_url, base_url)
            self._finish_slide(key, slide_id, version, status=SLIDE_STATUS_SUCCESS, image_url=image_url, error="", image_metadata=image_metadata)
        except Exception as exc:
            self._finish_slide(key, slide_id, version, status=SLIDE_STATUS_ERROR, error=str(exc) or "图片编辑失败")
        self._finalize_task_status(key)

    def _finish_slide(self, key: str, slide_id: str, version: int, *, status: str, image_url: str = "", error: str = "", image_metadata: dict[str, int] | None = None) -> None:
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            if self._task_cancel_requested_locked(task):
                return
            slide = self._find_slide_locked(task, slide_id)
            if slide is None or int(slide.get("version") or 0) != version:
                return
            if slide.get("status") == SLIDE_STATUS_STOPPED:
                return
            slide["status"] = status
            slide["error"] = error
            if image_url:
                slide["image_url"] = image_url
            if image_metadata:
                slide.update(image_metadata)
            slide["finished_at"] = _now_iso()
            task["updated_at"] = _now_iso()
            self._save_locked()

    def _finalize_task_status(self, key: str) -> None:
        with self._lock:
            task = self._tasks.get(key)
            if task is None:
                return
            self._finalize_task_status_locked(task)
            task["updated_at"] = _now_iso()
            self._save_locked()

    def _finalize_task_status_locked(self, task: dict[str, Any]) -> None:
        slides = self._slides_locked(task)
        if self._task_cancel_requested_locked(task) or any(slide.get("status") == SLIDE_STATUS_STOPPED for slide in slides):
            task["status"] = TASK_STATUS_STOPPED
            task["error"] = "任务已停止"
        elif any(slide.get("status") in UNFINISHED_SLIDE_STATUSES for slide in slides):
            task["status"] = TASK_STATUS_RUNNING
        elif slides and all(slide.get("status") == SLIDE_STATUS_SUCCESS for slide in slides):
            task["status"] = TASK_STATUS_SUCCESS
            task["error"] = ""
        elif any(slide.get("status") == SLIDE_STATUS_DRAFT for slide in slides):
            task["status"] = TASK_STATUS_DRAFT
            task["error"] = ""
        elif any(slide.get("status") == SLIDE_STATUS_ERROR for slide in slides):
            task["status"] = TASK_STATUS_ERROR
            task["error"] = "部分页面生成失败，请修改提示词后重试"
        else:
            task["status"] = TASK_STATUS_ERROR
            task["error"] = "没有可用页面"
        task["updated_at"] = _now_iso()
        if task.get("status") in {TASK_STATUS_SUCCESS, TASK_STATUS_ERROR, TASK_STATUS_STOPPED, TASK_STATUS_PACKAGED}:
            task["finished_at"] = task["updated_at"]

    @staticmethod
    def _invalidate_package_locked(task: dict[str, Any]) -> None:
        task.pop("pptx_path", None)
        task.pop("packaged_at", None)
        task.pop("pptx_build_version", None)

    def _package_path(self, owner: str, task_id: str) -> Path:
        return self.package_dir / _safe_file_stem(owner, "owner") / f"{_safe_file_stem(task_id, uuid.uuid4().hex)}.pptx"

    @staticmethod
    def _public_slide(slide: dict[str, Any]) -> dict[str, Any]:
        return {
            "slide_id": slide.get("slide_id"),
            "title": slide.get("title"),
            "layout_type": slide.get("layout_type") or MASTER_LAYOUT_SINGLE_COLUMN,
            "chapter_no": slide.get("chapter_no") or "",
            "chapter_title": slide.get("chapter_title") or "",
            "original_prompt": _strip_generated_chapter_instructions(_clean(slide.get("original_prompt"))),
            "current_prompt": _strip_generated_chapter_instructions(_clean(slide.get("current_prompt"))),
            "final_prompt": slide.get("final_prompt"),
            "image_url": slide.get("image_url"),
            "image_size": int(slide.get("image_size") or 0),
            "image_width": int(slide.get("image_width") or 0),
            "image_height": int(slide.get("image_height") or 0),
            "version": slide.get("version"),
            "status": slide.get("status"),
            "error": slide.get("error"),
            "started_at": slide.get("started_at") or "",
            "finished_at": slide.get("finished_at") or "",
            "reference_images": _normalize_reference_images(slide.get("reference_images")),
        }

    @staticmethod
    def _slides_locked(task: dict[str, Any]) -> list[dict[str, Any]]:
        slides = task.get("slides")
        return slides if isinstance(slides, list) else []

    @staticmethod
    def _task_cancel_requested_locked(task: dict[str, Any]) -> bool:
        return bool(task.get("cancel_requested")) or task.get("status") == TASK_STATUS_STOPPED

    @staticmethod
    def _find_slide_locked(task: dict[str, Any], slide_id: str) -> dict[str, Any] | None:
        for slide in PptTaskService._slides_locked(task):
            if _clean(slide.get("slide_id")) == slide_id:
                return slide
        return None

    @staticmethod
    def _slide_reference_image_locked(task: dict[str, Any], slide: dict[str, Any]) -> str:
        references = PptTaskService._slide_reference_images_locked(task, slide)
        return _clean(references[0].get("image_url")) if references else ""

    @staticmethod
    def _slide_reference_images_locked(task: dict[str, Any], slide: dict[str, Any]) -> list[dict[str, str]]:
        if "reference_images" in slide:
            return _normalize_reference_images(slide.get("reference_images"))
        if _clean(task.get("task_type"), TASK_TYPE_CONTENT) == TASK_TYPE_MASTER:
            layout_type = _normalize_master_layout(slide.get("layout_type") or slide.get("slide_id"))
            if layout_type == MASTER_BASE_LAYOUT:
                return []
            for master_slide in PptTaskService._slides_locked(task):
                if _normalize_master_layout(master_slide.get("layout_type") or master_slide.get("slide_id")) != MASTER_BASE_LAYOUT:
                    continue
                image_url = _clean(master_slide.get("image_url"))
                if image_url:
                    return [_reference_image_from_slide(master_slide)]
            raise ValueError("请先生成封面页母版后再生成其他母版")
        master_slides = task.get("master_slides")
        if not isinstance(master_slides, list) or not master_slides:
            return []
        layout_type = _normalize_master_layout(slide.get("layout_type"))
        for master_slide in master_slides:
            if not isinstance(master_slide, dict):
                continue
            if _normalize_master_layout(master_slide.get("layout_type") or master_slide.get("slide_id")) != layout_type:
                continue
            image_url = _clean(master_slide.get("image_url"))
            if image_url:
                return [_reference_image_from_slide(master_slide)]
        raise ValueError(f"未找到 {_master_layout_label(layout_type)} 对应的母版参考图")

    @staticmethod
    def _default_content_references(master_slides: list[dict[str, Any]], layout_type: object) -> list[dict[str, str]]:
        layout = _normalize_master_layout(layout_type)
        for master_slide in master_slides:
            if not isinstance(master_slide, dict):
                continue
            if _normalize_master_layout(master_slide.get("layout_type") or master_slide.get("slide_id")) != layout:
                continue
            if _clean(master_slide.get("image_url")):
                return [_reference_image_from_slide(master_slide)]
        return []

    def _validate_master_task_locked(
        self,
        task: dict[str, Any] | None,
        task_id: str,
        *,
        require_confirmed: bool,
    ) -> None:
        if task is None:
            raise PptTaskNotFoundError(task_id)
        if _clean(task.get("task_type"), TASK_TYPE_CONTENT) != TASK_TYPE_MASTER:
            raise ValueError("请选择母版任务")
        slides = self._slides_locked(task)
        by_layout = {_normalize_master_layout(slide.get("layout_type") or slide.get("slide_id")): slide for slide in slides}
        missing = [layout for layout in MASTER_LAYOUT_ORDER if layout not in by_layout]
        if missing:
            raise ValueError("母版任务缺少 8 类页面")
        unfinished = [slide for slide in by_layout.values() if slide.get("status") in UNFINISHED_SLIDE_STATUSES]
        if unfinished:
            raise ValueError("母版仍在生成，请稍后再试")
        failed_or_blank = [
            slide
            for layout, slide in by_layout.items()
            if layout in MASTER_LAYOUTS and (slide.get("status") != SLIDE_STATUS_SUCCESS or not _clean(slide.get("image_url")))
        ]
        if failed_or_blank:
            raise ValueError("8 张母版图片都生成成功后才能继续")
        if require_confirmed and not bool(task.get("master_confirmed")):
            raise ValueError("请先确认母版后再继续生成内容")

    def _master_snapshot_locked(self, owner: str, task_id: str, *, require_confirmed: bool) -> list[dict[str, Any]]:
        task = self._tasks.get(_task_key(owner, task_id))
        self._validate_master_task_locked(task, task_id, require_confirmed=require_confirmed)
        slides = self._slides_locked(task)
        by_layout = {_normalize_master_layout(slide.get("layout_type") or slide.get("slide_id")): slide for slide in slides}
        return self._public_master_slides([by_layout[layout] for layout in MASTER_LAYOUT_ORDER])

    def _public_task(self, task: dict[str, Any]) -> dict[str, Any]:
        slides: list[dict[str, Any]] = []
        for slide in self._slides_locked(task):
            public_slide = self._public_slide(slide)
            if "reference_images" not in slide:
                try:
                    public_slide["reference_images"] = self._slide_reference_images_locked(task, slide)
                except Exception:
                    public_slide["reference_images"] = []
            slides.append(public_slide)
        item = {
            "id": task.get("id"),
            "name": task.get("name") or task.get("id"),
            "task_type": task.get("task_type") or TASK_TYPE_CONTENT,
            "master_task_id": task.get("master_task_id") or "",
            "master_confirmed": bool(task.get("master_confirmed")),
            "status": task.get("status"),
            "slide_count": task.get("slide_count"),
            "design_concept": task.get("design_concept"),
            "global_style_prompt": task.get("global_style_prompt"),
            "chapters": task.get("chapters") if isinstance(task.get("chapters"), list) else [],
            "master_style_prompt": task.get("master_style_prompt") or "",
            "markdown": task.get("markdown") or "",
            "markdown_file_name": task.get("markdown_file_name") or "",
            "model": task.get("model"),
            "account_type": normalize_image_account_type(task.get("account_type")),
            "size": normalize_ppt_image_size(task.get("size")),
            "quality": normalize_image_quality(task.get("quality")),
            "concurrency": task.get("concurrency", DEFAULT_IMAGE_CONCURRENCY),
            "image_base_url": task.get("image_base_url"),
            "created_at": task.get("created_at"),
            "updated_at": task.get("updated_at"),
            "started_at": task.get("started_at") or "",
            "finished_at": task.get("finished_at") or "",
            "slides": slides,
            "pptx_ready": bool(
                _clean(task.get("pptx_path"))
                and _pptx_build_version(task.get("pptx_build_version")) == PPTX_BUILD_VERSION
                and Path(_clean(task.get("pptx_path"))).is_file()
            ),
        }
        master_slides = task.get("master_slides")
        if isinstance(master_slides, list) and master_slides:
            item["master_slides"] = self._public_master_slides(master_slides)
        if task.get("pptx_build_version"):
            item["pptx_build_version"] = task.get("pptx_build_version")
        if task.get("packaged_at"):
            item["packaged_at"] = task.get("packaged_at")
        if item["pptx_ready"]:
            item["download_url"] = f"/api/ppt/tasks/{task.get('id')}/download"
        if task.get("error"):
            item["error"] = task.get("error")
        return item

    def _public_master_slides(self, master_slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
        public_slides = [self._public_slide(slide) for slide in master_slides if isinstance(slide, dict)]
        base_slide = next(
            (
                slide
                for slide in public_slides
                if _normalize_master_layout(slide.get("layout_type") or slide.get("slide_id")) == MASTER_BASE_LAYOUT
                and _clean(slide.get("image_url"))
            ),
            None,
        )
        if not base_slide:
            return public_slides
        base_reference = _reference_image_from_slide(base_slide)
        for raw_slide, public_slide in zip([slide for slide in master_slides if isinstance(slide, dict)], public_slides):
            layout_type = _normalize_master_layout(public_slide.get("layout_type") or public_slide.get("slide_id"))
            if layout_type == MASTER_BASE_LAYOUT:
                public_slide["reference_images"] = _normalize_reference_images(raw_slide.get("reference_images"))
            elif "reference_images" not in raw_slide:
                public_slide["reference_images"] = [base_reference]
        return public_slides

    def _load_locked(self) -> dict[str, dict[str, Any]]:
        if not self.path.exists():
            return {}
        try:
            raw = json.loads(self.path.read_text(encoding="utf-8"))
        except Exception:
            return {}
        raw_items = raw.get("tasks") if isinstance(raw, dict) else raw
        if not isinstance(raw_items, list):
            return {}
        tasks: dict[str, dict[str, Any]] = {}
        for item in raw_items:
            if not isinstance(item, dict):
                continue
            task_id = _clean(item.get("id"))
            owner = _clean(item.get("owner_id"))
            if not task_id or not owner:
                continue
            slides: list[dict[str, Any]] = []
            for index, raw_slide in enumerate(item.get("slides") if isinstance(item.get("slides"), list) else [], start=1):
                if not isinstance(raw_slide, dict):
                    continue
                status = _clean(raw_slide.get("status"), SLIDE_STATUS_ERROR)
                if status not in SLIDE_STATUSES:
                    status = SLIDE_STATUS_ERROR
                slide_item = {
                    "slide_id": _clean(raw_slide.get("slide_id"), str(index)),
                    "title": _clean(raw_slide.get("title"), f"第 {index} 页"),
                    "layout_type": _normalize_master_layout(raw_slide.get("layout_type") or raw_slide.get("slide_id")),
                    "chapter_no": _clean(raw_slide.get("chapter_no")),
                    "chapter_title": _clean(raw_slide.get("chapter_title")),
                    "original_prompt": _clean(raw_slide.get("original_prompt")),
                    "current_prompt": _clean(raw_slide.get("current_prompt")),
                    "final_prompt": _clean(raw_slide.get("final_prompt")),
                    "image_url": _clean(raw_slide.get("image_url")),
                    "image_size": int(raw_slide.get("image_size") or 0),
                    "image_width": int(raw_slide.get("image_width") or 0),
                    "image_height": int(raw_slide.get("image_height") or 0),
                    "version": int(raw_slide.get("version") or 1),
                    "status": status,
                    "error": _clean(raw_slide.get("error")),
                    "started_at": _clean(raw_slide.get("started_at")),
                    "finished_at": _clean(raw_slide.get("finished_at")),
                }
                if "reference_images" in raw_slide:
                    slide_item["reference_images"] = _normalize_reference_images(raw_slide.get("reference_images"))
                slides.append(slide_item)
            status = _clean(item.get("status"), TASK_STATUS_ERROR)
            if status not in TASK_STATUSES:
                status = TASK_STATUS_ERROR
            task_type = _clean(item.get("task_type"), TASK_TYPE_CONTENT)
            if task_type not in TASK_TYPES:
                task_type = TASK_TYPE_CONTENT
            try:
                concurrency = normalize_image_concurrency(item.get("concurrency"))
            except ValueError:
                concurrency = DEFAULT_IMAGE_CONCURRENCY
            master_slides: list[dict[str, Any]] = []
            for index, raw_slide in enumerate(item.get("master_slides") if isinstance(item.get("master_slides"), list) else [], start=1):
                if not isinstance(raw_slide, dict):
                    continue
                master_status = _clean(raw_slide.get("status"), SLIDE_STATUS_ERROR)
                if master_status not in SLIDE_STATUSES:
                    master_status = SLIDE_STATUS_ERROR
                layout_type = _normalize_master_layout(raw_slide.get("layout_type") or raw_slide.get("slide_id"))
                master_slide_item = {
                    "slide_id": _clean(raw_slide.get("slide_id"), layout_type or str(index)),
                    "title": _clean(raw_slide.get("title"), MASTER_LAYOUT_LABELS.get(layout_type, f"母版 {index}")),
                    "layout_type": layout_type,
                    "chapter_no": _clean(raw_slide.get("chapter_no")),
                    "chapter_title": _clean(raw_slide.get("chapter_title")),
                    "original_prompt": _clean(raw_slide.get("original_prompt")),
                    "current_prompt": _clean(raw_slide.get("current_prompt")),
                    "final_prompt": _clean(raw_slide.get("final_prompt")),
                    "image_url": _clean(raw_slide.get("image_url")),
                    "image_size": int(raw_slide.get("image_size") or 0),
                    "image_width": int(raw_slide.get("image_width") or 0),
                    "image_height": int(raw_slide.get("image_height") or 0),
                    "version": int(raw_slide.get("version") or 1),
                    "status": master_status,
                    "error": _clean(raw_slide.get("error")),
                    "started_at": _clean(raw_slide.get("started_at")),
                    "finished_at": _clean(raw_slide.get("finished_at")),
                }
                if "reference_images" in raw_slide:
                    master_slide_item["reference_images"] = _normalize_reference_images(raw_slide.get("reference_images"))
                master_slides.append(master_slide_item)
            task = {
                "id": task_id,
                "name": _clean(item.get("name"), task_id),
                "owner_id": owner,
                "task_type": task_type,
                "master_confirmed": bool(item.get("master_confirmed")),
                "status": status,
                "model": _clean(item.get("model"), DEFAULT_IMAGE_MODEL),
                "account_type": normalize_image_account_type(item.get("account_type")),
                "size": normalize_ppt_image_size(item.get("size")),
                "quality": normalize_image_quality(item.get("quality")),
                "concurrency": concurrency,
                "image_base_url": _clean(item.get("image_base_url")),
                "image_api_key": _clean(item.get("image_api_key")),
                "slide_count": int(item.get("slide_count") or len(slides) or DEFAULT_SLIDE_COUNT),
                "design_concept": _clean(item.get("design_concept")),
                "global_style_prompt": _clean(item.get("global_style_prompt"), GLOBAL_STYLE_PROMPT),
                "chapters": _normalize_plan_chapters(item.get("chapters")),
                "master_style_prompt": _clean(item.get("master_style_prompt")),
                "markdown": _clean(item.get("markdown")),
                "markdown_file_name": _clean(item.get("markdown_file_name")),
                "slides": slides,
                "created_at": _clean(item.get("created_at"), _now_iso()),
                "updated_at": _clean(item.get("updated_at"), _clean(item.get("created_at"), _now_iso())),
                "started_at": _clean(item.get("started_at")),
                "finished_at": _clean(item.get("finished_at")),
                "error": _clean(item.get("error")),
            }
            if bool(item.get("cancel_requested")) and status == TASK_STATUS_STOPPED:
                task["cancel_requested"] = True
            if _clean(item.get("master_task_id")):
                task["master_task_id"] = _clean(item.get("master_task_id"))
            if master_slides:
                task["master_slides"] = master_slides
            if _clean(item.get("pptx_path")):
                task["pptx_path"] = _clean(item.get("pptx_path"))
            if _clean(item.get("pptx_build_version")):
                task["pptx_build_version"] = _pptx_build_version(item.get("pptx_build_version"))
            if _clean(item.get("packaged_at")):
                task["packaged_at"] = _clean(item.get("packaged_at"))
            tasks[_task_key(owner, task_id)] = task
        return tasks

    def _save_locked(self) -> None:
        items = sorted(self._tasks.values(), key=lambda item: str(item.get("updated_at") or ""), reverse=True)
        tmp_path = self.path.with_suffix(self.path.suffix + ".tmp")
        tmp_path.write_text(json.dumps({"tasks": items}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        tmp_path.replace(self.path)

    def _recover_unfinished_locked(self) -> bool:
        changed = False
        for task in self._tasks.values():
            if task.get("status") in UNFINISHED_TASK_STATUSES:
                task["status"] = TASK_STATUS_ERROR
                task["error"] = "服务已重启，未完成的 PPT 任务已中断"
                task["updated_at"] = _now_iso()
                changed = True
            for slide in self._slides_locked(task):
                if slide.get("status") in UNFINISHED_SLIDE_STATUSES:
                    slide["status"] = SLIDE_STATUS_ERROR
                    slide["error"] = "服务已重启，未完成的页面任务已中断"
                    changed = True
        return changed


ppt_task_service = PptTaskService(DATA_DIR / "ppt_tasks.json")
